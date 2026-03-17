import os
from pptx import Presentation

def extract_pptx_html(file_path):
    try:
        prs = Presentation(file_path)
        html = "<div style='background: #f0f0f0; padding: 20px; font-family: sans-serif;'>"
        for i, slide in enumerate(prs.slides):
            html += f"<div style='background: white; border: 1px solid #ccc; margin-bottom: 20px; padding: 40px; aspect-ratio: 16/9; position: relative;'><h3 style='color: #888; font-size: 14px;'>Slide {i+1}</h3>"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br>')
                    html += f"<p>{text}</p>"
            html += "</div>"
        html += "</div>"
        return html
    except Exception as e:
        return str(e)

if __name__ == "__main__":
    path = "uploads/GNN.pptx"
    if os.path.exists(path):
        result = extract_pptx_html(path)
        print(f"Extraction successful: {len(result)} chars")
        # print(result[:500])
    else:
        print("File not found")
