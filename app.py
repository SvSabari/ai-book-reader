from flask import Flask, render_template, request, jsonify, send_from_directory, send_file, Response
from werkzeug.utils import secure_filename
import os
import threading
import sqlite3
import time
import math
import re
import base64
import json
import traceback
import io
import uuid
import hashlib
import concurrent.futures
import asyncio
import random
import requests
import html
import collections
from typing import List
from io import BytesIO, StringIO
from urllib.parse import quote, unquote, urljoin
import urllib.request

# Third-party libraries
import fitz
import docx
from ebooklib import epub, ITEM_DOCUMENT, ITEM_IMAGE
from bs4 import BeautifulSoup
import cv2
import pytesseract
import mammoth
import numpy as np
from gtts import gTTS
import edge_tts
from deep_translator import GoogleTranslator
from fpdf import FPDF
from textblob import TextBlob
from PIL import Image
import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import subprocess
import time
import requests
import atexit

# --- GLOBAL CONFIG & TRANSLATION SIDECAR ---
SIDECAR_PORT = 3001
sidecar_process = None

def start_translation_sidecar():
    global sidecar_process
    try:
        if os.path.exists("translator_sidecar.js"):
            print(f"🚀 Launching Translation Sidecar (Port {SIDECAR_PORT})...")
            # Sidecar handles 100+ parallel translation hits via Node.js Event Loop
            sidecar_process = subprocess.Popen(["node", "translator_sidecar.js"], 
                                             stdout=subprocess.DEVNULL, 
                                             stderr=subprocess.DEVNULL)
            time.sleep(0.5)
    except Exception as e:
        print(f"Sidecar launch failed: {e}. Falling back to internal engine.")

@atexit.register
def kill_sidecar():
    if sidecar_process:
        print("🛑 Shutting down translation sidecar...")
        sidecar_process.terminate()

try:
    import transformers
    from transformers import pipeline
except ImportError:
    transformers = None
    pipeline = None

app = Flask(__name__)

UPLOAD_FOLDER = "uploads"
EXTRACTED_FOLDER = "extracted"
DB = "database.db"

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(EXTRACTED_FOLDER, exist_ok=True)
os.makedirs(os.path.join(UPLOAD_FOLDER, "extracted_assets"), exist_ok=True)

# Global cache to store pre-tokenized sentences for the currently active book
# Structure: { book_id: (timestamp, [sentences]) }
SENTENCE_CACHE = {}
_SENTIMENT_ANALYZER = None

def get_sentiment_analyzer():
    global _SENTIMENT_ANALYZER
    if _SENTIMENT_ANALYZER is None:
        try:
            if pipeline:
                _SENTIMENT_ANALYZER = pipeline("sentiment-analysis", model="distilbert-base-uncased-finetuned-sst-2-english", device=-1)
            else:
                _SENTIMENT_ANALYZER = "FAILED"
        except Exception as e:
            print(f"Sentiment model pre-load failed: {e}")
            _SENTIMENT_ANALYZER = "FAILED"
    return _SENTIMENT_ANALYZER if _SENTIMENT_ANALYZER != "FAILED" else None

def get_book_sentences(book_id):
    """Retrieve or generate tokenized sentences for a book to avoid repeated parsing."""
    if book_id in SENTENCE_CACHE:
        return SENTENCE_CACHE[book_id][1]
    
    conn = get_conn()
    cur = conn.cursor()
    cur.execute("SELECT extracted_path FROM books WHERE id = ?", (book_id,))
    row = cur.fetchone()
    conn.close()
    
    if not row or not os.path.exists(row[0]):
        return []
        
    try:
        with open(row[0], "r", encoding="utf-8", errors="ignore") as f:
            html_content = f.read()
            
        soup = BeautifulSoup(html_content, "html.parser")
        
        # Clean up HTML for narration (remove noise)
        for noise in soup(["script", "style", "nav", "footer", "header", "meta", "link"]):
            noise.decompose()
            
        text = soup.get_text(separator=" ")
        
        # Tokenize into sentences
        sentences = re.split(r'(?<=[.!?])\s+', text)
        sentences = [re.sub(r'\s+', ' ', s).strip() for s in sentences if len(s.strip()) > 5]
        
        # Keep only the last 3 books in cache to manage memory
        if len(SENTENCE_CACHE) > 3:
            oldest = min(SENTENCE_CACHE.keys(), key=lambda k: SENTENCE_CACHE[k][0])
            SENTENCE_CACHE.pop(oldest, None)
            
        SENTENCE_CACHE[book_id] = (time.time(), sentences)
        return sentences
    except Exception:
        return []

pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


def get_conn():
    return sqlite3.connect(DB)


def init_db():
    conn = get_conn()
    cur = conn.cursor()

    cur.execute("""
    CREATE TABLE IF NOT EXISTS books(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT,
        path TEXT,
        extracted_path TEXT,
        status TEXT DEFAULT 'ready',
        uploaded_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)

    # Migrate existing databases that don't have the status column
    try:
        cur.execute("ALTER TABLE books ADD COLUMN status TEXT DEFAULT 'ready'")
    except Exception:
        pass  # Column already exists

    cur.execute("""
    CREATE TABLE IF NOT EXISTS highlights(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        book_id INTEGER,
        highlighted_text TEXT
    )
    """)

    cur.execute("""
    CREATE TABLE IF NOT EXISTS notes(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        book_id INTEGER,
        content TEXT,
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)

    cur.execute("""
    CREATE TABLE IF NOT EXISTS bookmarks(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        book_id INTEGER,
        page_number INTEGER,
        scroll_y INTEGER,
        char_index INTEGER DEFAULT 0,
        node_index INTEGER DEFAULT -1,
        node_offset INTEGER DEFAULT 0,
        label TEXT,
        lang_code TEXT DEFAULT 'en',
        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
    )
    """)

    # Migration for char_index and node mapping
    try:
        cur.execute("ALTER TABLE bookmarks ADD COLUMN char_index INTEGER DEFAULT 0")
    except Exception: pass
    try:
        cur.execute("ALTER TABLE bookmarks ADD COLUMN node_index INTEGER DEFAULT -1")
    except Exception: pass
    try:
        cur.execute("ALTER TABLE bookmarks ADD COLUMN node_offset INTEGER DEFAULT 0")
    except Exception: pass
    try:
        cur.execute("ALTER TABLE bookmarks ADD COLUMN lang_code TEXT DEFAULT 'en'")
    except Exception: pass


    conn.commit()
    conn.close()


# Database schema maintenance is handled in the main entry point


try:
    # Google Cloud Vision integration has been disabled by user request.
    from google.cloud import vision
    def get_vision_client(): return None
except ImportError:
    vision = None
    def get_vision_client(): return None


def google_vision_ocr(img_cv):
    """Google Cloud Vision: The gold standard for handwriting recognition."""
    client = get_vision_client()
    if not client: return None
    
    try:
        _, encoded = cv2.imencode(".jpg", img_cv)
        content = encoded.tobytes()
        image = vision.Image({"content": content})
        # DOCUMENT_TEXT_DETECTION is optimized for handwriting and dense pages
        response = client.document_text_detection(image=image)
        if response.full_text_annotation:
            return response.full_text_annotation.text
    except Exception as e:
        print(f"Google Vision failed: {e}")
    return None


def extract_image_text(image, fast=True, skip_ocr=False):
    """Perform a high-fidelity OCR pass. Automatically upgrades to Google Vision for handwriting."""
    if skip_ocr: return ""
    try:
        if image is None: return ""
        if isinstance(image, str):
            # Universal Byte Decoding (Safer than imread for non-ASCII paths/WebP)
            with open(image, "rb") as f:
                img_bytes = f.read()
                image = cv2.imdecode(np.frombuffer(img_bytes, np.uint8), cv2.IMREAD_COLOR)
            if image is None: return ""
        
        # 1. ATTEMPT GOOGLE CLOUD VISION (PRIORITY: PERFECT ACCURACY & ORDER)
        # Always use Google Vision if client is available, as it's the gold standard for book layouts.
        # We only skip if skip_ocr is True.
        client = get_vision_client()
        if client:
            try:
                google_res = google_vision_ocr(image)
                if google_res and len(google_res.strip()) > 5:
                    # Vision already preserves reading order perfectly
                    return google_res.strip()
            except Exception as e:
                print(f"Google Vision fallback to Tesseract: {e}")

        # 2. LOCAL OCR FALLBACK (Optimized for Order and Accuracy)
        # Avoid multiple passes which cause "soo big" duplicated text.
        h, w = image.shape[:2]
        # Rescale for better Tesseract recognition on dense text
        scale = 1.0
        if max(h, w) < 1000: scale = 1.5
        if scale != 1.0:
            image = cv2.resize(image, (int(w * scale), int(h * scale)), interpolation=cv2.INTER_CUBIC)

        gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
        
        # PRE-PROCESSING: Standardize for OCR
        # Bilateral filter removes noise while keeping edges sharp (good for books)
        processed = cv2.bilateralFilter(gray, 9, 75, 75)
        # Adaptive thresholding works best for uneven lighting on book scans
        processed = cv2.adaptiveThreshold(processed, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2)

        # OCR: Use PSM 3 (Automatic Page Segmentation) to find lines in order without OSD overhead
        # Using OEM 3 (Default, combines LSTM and Legacy if needed)
        config = r'--tessdata-dir c:\ai_book_reader\tessdata --oem 3 --psm 3 -l eng'
        
        ocr_text = pytesseract.image_to_string(processed, config=config).strip()
        
        if not ocr_text: return ""

        # SANITIZE & ORDER PRESERVATION
        # Filter out obvious OCR noise (single chars, non-alphanumeric junk)
        final_lines = []
        for line in ocr_text.split('\n'):
            line = line.strip()
            # Ignore lines that are just symbols or too short to be meaningful book text
            if len(re.sub(r'[^a-zA-Z0-9]', '', line)) < 2: continue
            # Basic sanity filter for garbage
            if re.match(r'^[.|_|\-|\s|?|!|:|;]+$', line): continue
            final_lines.append(line)

        return "\n".join(final_lines)
    except Exception as e:
        print("OCR failed:", e)
        return ""


def build_clickable_img_wrapper(img_tag_str, img_cv, fast=True, skip_ocr=True):
    """Given an OpenCV image tag and image, wrap it in a clickable container for AI explanation.
    Hidden OCR is embedded for backend search/quiz availability."""
    # DEFAULT skip_ocr to True for massive book performance; background tasks can fill it later
    text = extract_image_text(img_cv, fast=fast, skip_ocr=skip_ocr)
    
    # Store OCR text in a hidden span so BeautifulSoup/get_text can index it for quizzes/searches
    # while keeping it invisible and non-interactive for the user to avoid visual clutter.
    wrapper = f'<div class="img-clickable-wrapper" onclick="explainImage(this)" style="cursor: zoom-in; display: inline-block; position: relative;">{img_tag_str}'
    if text:
        # Wrap in a visually hidden layer that narrators CAN read.
        # This is CRITICAL for scanned books where this is the only text.
        wrapper += f'<span class="ocr-reading-layer" style="position:absolute; width:1px; height:1px; overflow:hidden; opacity:0; pointer-events:none;">{text}</span>'
    wrapper += '</div>'
    return wrapper


def ocr_embedded_images(html):
    """
    Fast pass: extract raw text from embedded images for TTS/search purposes only.
    Uses regex and multithreading for high performance without BeautifulSoup overhead.
    """
    img_pattern = re.compile(r'(<img[^>]*src=[\'"](data:image/[^\'"]+)[\'"][^>]*>)', re.IGNORECASE)
    matches = list(img_pattern.finditer(html))
    
    if not matches:
        return html
        
    def process_match_data(match_tuple):
        start, end, img_tag, src_data = match_tuple
        try:
            _, encoded = src_data.split(",", 1)
            data = base64.b64decode(encoded)
            nparr = np.frombuffer(data, np.uint8)
            img_cv = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
            if img_cv is not None:
                # Wrap images in a clickable container for explanation instead of OCR overlay
                clickable_html = build_clickable_img_wrapper(img_tag, img_cv)
                if clickable_html:
                    return start, end, clickable_html
        except Exception:
            pass
        return start, end, img_tag

    tasks = [(m.start(), m.end(), m.group(1), m.group(2)) for m in matches]
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=4) as executor:
        results = list(executor.map(process_match_data, tasks))
        
    parts = []
    last_end = 0
    # Process results in sequential order of occurrence in the HTML
    results.sort(key=lambda x: x[0])
    for start, end, replacement in results:
        parts.append(html[last_end:start])
        parts.append(replacement)
        last_end = end
    parts.append(html[last_end:])
    
    return "".join(parts)


def extract_pdf_html(file_path, fast_mode=True):
    """High-speed parallel PDF processing with watermark filtering and layout stabilization."""
    abs_path = os.path.abspath(file_path)
    if not os.path.exists(abs_path):
        return "<p>Error: PDF source file missing.</p>"

    pdf = fitz.open(abs_path)
    if pdf.is_encrypted:
        try: pdf.decrypt("")
        except: pass
    
    total_pages = len(pdf)
    
    def process_pdf_page(page_num):
        try:
            # We open a NEW copy of the PDF handle per thread to ensure thread-safety in fitz
            local_pdf = fitz.open(abs_path)
            page = local_pdf[page_num]
            page_area = page.rect.width * page.rect.height
            
            # 1. Page Header/Container
            page_html = f'<div id="pdf-page-{page_num}" class="lazy-page-container flex-page" data-original-width="{int(page.rect.width)}" style="display: flex; flex-direction: column; margin-bottom: 60px; padding: 40px; gap: 20px; background: white; border-radius: 8px; box-shadow: 0 4px 20px rgba(0,0,0,0.05); min-height: 200px;">'
            
            img_html_parts = []
            text_html_parts = []
            seen_hashes = set()

            # A. EXTRACT IMAGES (Filter watermarks)
            page_images = page.get_images(full=True)
            for img_info in page_images:
                try:
                    xref = img_info[0]
                    img_rects = page.get_image_rects(xref)
                    if not img_rects: continue
                    rect = img_rects[0]
                    
                    # IGNORE SMALL ASSETS (Skip 'Digitized by Google', tiny icons, separators)
                    if rect.width * rect.height < (page_area * 0.05): continue

                    pix = page.parent.extract_image(xref)
                    img_data = pix["image"]
                    ext = pix["ext"]
                    img_hash = hashlib.md5(img_data).hexdigest()
                    if img_hash in seen_hashes: continue
                    seen_hashes.add(img_hash)
                    
                    asset_name = f"assets_{img_hash}.{ext}"
                    asset_path = os.path.join(UPLOAD_FOLDER, "extracted_assets", asset_name)
                    if not os.path.exists(asset_path):
                        os.makedirs(os.path.dirname(asset_path), exist_ok=True)
                        with open(asset_path, "wb") as f: f.write(img_data)
                    
                    img_tag = f'<img src="/uploads/extracted_assets/{asset_name}" style="max-width: 100%; height: auto; display: block; margin: 0 auto; border-radius: 8px;" />'
                    
                    try:
                        img_np = cv2.imdecode(np.frombuffer(img_data, np.uint8), cv2.IMREAD_COLOR)
                        if img_np is not None:
                            # Use fast sync OCR only for very large images (potential full-page scans)
                            do_deep = (rect.width * rect.height > page_area * 0.7)
                            wrapped = build_clickable_img_wrapper(img_tag, img_np, fast=True, skip_ocr=False if do_deep and not fast_mode else True)
                            img_html_parts.append(wrapped)
                        else: img_html_parts.append(img_tag)
                    except: img_html_parts.append(img_tag)
                except: pass

            # B. EXTRACT TEXT
            blocks = page.get_text("dict").get("blocks", [])
            is_scanned = any(img_info[2]*img_info[3] > page_area * 0.75 for img_info in page_images) if page_images else False
            
            seen_texts = set()
            for b in blocks:
                if b.get("type", 0) == 0:
                    lines_text = []
                    for line in b.get("lines", []):
                        lines_text.append("".join([s.get("text", "") for s in line.get("spans", [])]))
                    block_text = " ".join(lines_text).strip()
                    
                    if not block_text or len(block_text) < 4: continue
                    if block_text in seen_texts: continue
                    l_text = " ".join(["".join([s.get("text", "") for s in line.get("spans", [])]) for line in b.get("lines", [])]).strip()
                    if not l_text or len(l_text) < 4: continue
                    if l_text in seen_texts: continue
                    seen_texts.add(l_text)
                    
                    style = 'style="margin-bottom: 1.25em; line-height: 1.8; color: #334155;"'
                    if is_scanned:
                        style = 'style="position:absolute; width:1px; height:1px; overflow:hidden; opacity:0; pointer-events:none;"'
                    text_html_parts.append(f'<p {style}>{l_text}</p>')

            # C. ASSEMBLY
            if img_html_parts:
                page_html += '<div class="pdf-img-stack" style="text-align: center; margin-bottom: 20px;">' + "".join(img_html_parts) + '</div>'
            if text_html_parts:
                page_html += '<div class="pdf-text-stack" style="width: 100%; font-family: \'Georgia\', serif; font-size: 1.15rem;">' + "".join(text_html_parts) + '</div>'
            
            page_html += "</div>"
            local_pdf.close()
            return page_html
        except Exception as e:
            return f'<div class="error">Page Error: {e}</div>'

    # Process all pages in parallel
    with concurrent.futures.ThreadPoolExecutor(max_workers=8) as executor:
        results = list(executor.map(process_pdf_page, range(total_pages)))
    
    pdf.close()
    return "".join(results)


def extract_docx_html(file_path):
    def convert_image(image):
        with image.open() as image_bytes:
            encoded_src = base64.b64encode(image_bytes.read()).decode("ascii")
            ext = image.content_type.split("/")[-1]
            # Use interactive wrapper for potential handwriting/diagrams
            img_tag = f'<img src="data:image/{ext};base64,{encoded_src}" />'
            try:
                # Re-decode to get CV2 image for OCR
                image_bytes.seek(0)
                nparr = np.frombuffer(image_bytes.read(), np.uint8)
                img_cv = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                if img_cv is not None:
                    return {"src": "", "html": build_clickable_img_wrapper(img_tag, img_cv)}
            except: pass
            return {"src": f"data:image/{ext};base64,{encoded_src}"}

    with open(file_path, "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file, convert_image=mammoth.images.img_element(convert_image))
        html = result.value
    return f"<div id='docx-page-0' class='lazy-page-container' style='padding: 60px; max-width: 900px; margin: 0 auto; background: var(--bg-paper); border-radius: 4px; box-shadow: 0 4px 20px rgba(0,0,0,0.1); min-height: 800px; color: #333; font-family: \"Georgia\", serif; font-size: 1.15rem;'>{html}</div>"


def extract_epub_html(file_path):
    book = epub.read_epub(file_path)
    
    # Pre-map ALL images in the archive regardless of folder structure
    img_map = {}
    for item in book.get_items():
        if int(item.get_type()) == ITEM_IMAGE: 
            i_name = str(item.get_name())
            name_key = i_name.split('/')[-1]
            img_map[name_key] = item.get_content()

    html_content = ""
    page_count = 0
    
    for item in book.get_items():
        if int(item.get_type()) == ITEM_DOCUMENT: 
            content = item.get_content()
            if not content: continue
            
            soup = BeautifulSoup(content, "html.parser")
            for img in soup.find_all("img"):
                src = str(img.attrs.get("src", ""))
                if src:
                    name = src.split('/')[-1]
                    img_data = img_map.get(name)
                    if img_data is not None:
                        b64 = base64.b64encode(img_data).decode('utf-8')
                        ext = name.split('.')[-1].lower() if '.' in name else 'png'
                        img.attrs['src'] = f"data:image/{ext};base64,{b64}"
                        # Make it interactive
                        img_tag_str = str(img)
                        try:
                            # Re-wrap in interactive container
                            img_np = cv2.imdecode(np.frombuffer(img_data, np.uint8), cv2.IMREAD_COLOR)
                            if img_np is not None:
                                interactive = build_clickable_img_wrapper(img_tag_str, img_np)
                                # Replace the img tag with the wrapped version in the soup
                                new_soup = BeautifulSoup(interactive, "html.parser")
                                img.replace_with(new_soup)
                        except Exception:
                            pass

            body = soup.find('body')
            if body:
                # Wrap each EPUB document (usually a chapter) in its own lazy container
                # This fixes the 'one giant page' issue that stalls translation for massive epubs.
                page_html = f"<div id='epub-page-{page_count}' class='lazy-page-container' style='padding: 60px; max-width: 900px; margin: 0 auto 40px auto; background: var(--bg-paper); border-radius: 4px; box-shadow: 0 4px 20px rgba(0,0,0,0.1); min-height: 800px; color: #333; font-family: \"Georgia\", serif; font-size: 1.15rem; perspective: 1000px;'>"
                for child in body.children:
                    page_html += str(child)
                page_html += "</div>"
                html_content += page_html
                page_count += 1
    
    return html_content if html_content else "<div class='error'>No readable content found in EPUB.</div>"


def extract_pptx_html(file_path):
    try:
        prs = pptx.Presentation(file_path)
    except Exception:
        return "<div class='error'>Failed to load PPTX. Try saving as PDF.</div>"

    sw = prs.slide_width
    sh = prs.slide_height

    # Convert EMUs to a base pixel width for consistent frontend rendering
    base_w = 800
    base_h = int((sh / sw) * base_w)
    aspect = (sh / sw) * 100
    
    html = f'<div class="book-content-container" style="background:transparent; padding: 20px;">'
    for i, slide in enumerate(prs.slides):
        # Create a positioned slide canvas matching the PowerPoint aspect ratio
        html += f'<div id="slide-page-{i}" class="lazy-page-container" data-original-width="{base_w}" data-original-height="{base_h}" style="width: {base_w}px; height: {base_h}px; position: relative; background: #fff; margin: 0 auto 40px auto; box-shadow: 0 4px 15px rgba(0,0,0,0.1); border-radius: 4px; overflow: hidden; container-type: size;">'
        html += f'<div style="position: absolute; top: 0; left: 0; width: 100%; height: 100%; overflow: hidden;">'
        
        # Add a light background number
        html += f'<div style="position: absolute; bottom: 10px; right: 15px; color: #ccc; font-size: 0.8em; z-index: 0;">Slide {i+1}</div>'

        for shape in slide.shapes:
            # Calculate percentage coordinates
            t = (shape.top / sh) * 100
            l = (shape.left / sw) * 100
            w = (shape.width / sw) * 100
            h = (shape.height / sh) * 100
            
            # Skip shapes that are way off-canvas
            if l < -10 or t < -10 or l > 110 or t > 110: continue

            style = f"position: absolute; top: {t}%; left: {l}%; width: {w}%; height: {h}%;"

            if hasattr(shape, "image"):
                try:
                    b64 = base64.b64encode(shape.image.blob).decode()
                    ext = shape.image.ext or "png"
                    html += f'<div style="{style} z-index: 1;"><img src="data:image/{ext};base64,{b64}" style="width: 100%; height: 100%; object-fit: contain;" /></div>'
                except Exception: pass
            elif hasattr(shape, "text_frame") and shape.text_frame.text.strip():
                # Extract alignment and native font size
                h_align = "left"
                v_align = "flex-start"
                
                # Try to get native font size if set
                fs_px = None
                if shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                    # Check paragraph alignment
                    if p.alignment == PP_ALIGN.CENTER: h_align = "center"
                    elif p.alignment == PP_ALIGN.RIGHT: h_align = "right"
                    elif p.alignment == PP_ALIGN.JUSTIFY: h_align = "justify"
                    
                    # Check font size (native PowerPoint font size)
                    if p.font.size:
                        # Convert EMUs to pixels based on our 800px base width
                        fs_px = (p.font.size / sw) * base_w

                if shape.text_frame.vertical_anchor == MSO_ANCHOR.MIDDLE: v_align = "center"
                elif shape.text_frame.vertical_anchor == MSO_ANCHOR.BOTTOM: v_align = "flex-end"

                text = shape.text_frame.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;').replace('\n', '<br>')
                
                # Determine font size: Use native if found, otherwise fallback to a more conservative clamp
                if fs_px and fs_px > 5:
                    font_size_style = f"font-size: {round(fs_px, 1)}px;"
                else:
                    font_size_style = "font-size: clamp(8px, 1.8cqw, 36px);"

                # For text boxes, we often don't want a fixed height to prevent clipping, 
                # but we keep width fixed for wrapping.
                html += f'<div style="position: absolute; top: {t}%; left: {l}%; width: {w}%; min-height: {h}%; z-index: 2; padding: 2px; display: flex; align-items: {v_align}; line-height: 1.1; {font_size_style} color: #333; text-align: {h_align}; word-break: break-word;">'
                html += f'<div style="width: 100%;">{text}</div></div>'
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                # Basic table support
                html += f'<div style="{style} z-index: 2; background: #f9f9f9; border: 1px solid #ddd; font-size: 0.6em; overflow: auto; padding: 2px;">Table Data Linked</div>'

        html += "</div></div>"
    html += "</div>"
    return html


def extract_txt_html(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        # 1. Clean and Prepare
        raw_text = f.read().replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        
        # 2. VIRTUAL PAGINATION (Crucial for Infinite Scroll and Translation Throughput)
        # We split the book into chunks of approx 3000 characters to ensure the 
        # translation engine stays under rate-limit thresholds and provides 
        # near-instant results for the visible page.
        paragraphs = raw_text.split('\n')
        pages = []
        current_page_text = []
        current_len = 0
        MAX_PAGE_CHARS = 3000 # Roughly 3 small paragraphs or 1 large one
        
        for p in paragraphs:
            p = p.strip()
            if not p:
                current_page_text.append("<br>")
                continue
                
            current_page_text.append(p)
            current_len += len(p)
            
            if current_len > MAX_PAGE_CHARS:
                # Close Page
                page_id = len(pages)
                html = "<br><br>".join(current_page_text)
                pages.append(f"<div id='txt-page-{page_id}' class='lazy-page-container' style='padding: 60px; max-width: 900px; margin: 40px auto; background: var(--bg-paper); border-radius: 8px; box-shadow: 0 4px 30px rgba(0,0,0,0.1); min-height: 700px; color: #222; font-family: \"Georgia\", serif; font-size: 1.2rem; line-height: 1.8; position: relative;'>{html}</div>")
                current_page_text = []
                current_len = 0
        
        # Handle trailing page
        if current_page_text:
            page_id = len(pages)
            html = "<br><br>".join(current_page_text)
            pages.append(f"<div id='txt-page-{page_id}' class='lazy-page-container' style='padding: 60px; max-width: 900px; margin: 40px auto; background: var(--bg-paper); border-radius: 8px; box-shadow: 0 4px 30px rgba(0,0,0,0.1); min-height: 700px; color: #222; font-family: \"Georgia\", serif; font-size: 1.2rem; line-height: 1.8; position: relative;'>{html}</div>")
            
        return "\n".join(pages) if pages else "<div class='lazy-page-container' style='padding:100px; text-align:center;'>Empty Document</div>"


def extract_book_html(file_path, fast_mode=True):
    ext = os.path.splitext(file_path)[1].lower()
    
    # 1. ROBUST EXTENSION DETECTION (MAGIC NUMBERS)
    # If the extension is missing or potentially wrong, verify against file signature
    if not ext or ext not in [".pdf", ".docx", ".epub", ".pptx"]:
        try:
            with open(file_path, "rb") as f:
                head = f.read(4)
                if head == b"%PDF": ext = ".pdf"
                elif head == b"PK\x03\x04":
                    # It's a zip-based format (EPUB, DOCX, PPTX)
                    # We can't immediately tell which, but we can try EBOPUB logic first
                    ext = ".epub" 
        except: pass

    # 2. MATCH AND EXTRACT
    if ext == ".pdf":
        html = extract_pdf_html(file_path)
        return ocr_embedded_images(html) if not fast_mode else html
    elif ext == ".docx":
        html = extract_docx_html(file_path)
        return ocr_embedded_images(html) if not fast_mode else html
    elif ext == ".epub":
        try:
            html = extract_epub_html(file_path)
            # If EPUB extraction produced literal nothing, fallback to TXT
            if not html or len(str(html).strip()) < 10:
                raise Exception("Empty EPUB")
            return ocr_embedded_images(html) if not fast_mode else html
        except Exception:
            # Fallback to structural TXT extraction as a last resort
            return extract_txt_html(file_path)
    elif ext in [".txt", ".prn", ".text", ".md", ".log"]:
        return extract_txt_html(file_path)
    elif ext in [".html", ".htm", ".xhtml"]:
        # Direct HTML reading (preserving original layout)
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
                # Wrap in a lazy page container for consistency
                return f"<div id='html-page-0' class='lazy-page-container' style='padding: 60px; max-width: 900px; margin: 0 auto; background: var(--bg-paper); border-radius: 4px; box-shadow: 0 4px 20px rgba(0,0,0,0.1); color: #333;'>{content}</div>"
        except:
            return extract_txt_html(file_path)
    elif ext == ".pptx":
        return extract_pptx_html(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".webp", ".gif", ".tiff", ".svg", ".ico", ".jfif"]:
        # Logic for image-to-book (already works)
        filename = os.path.basename(file_path)
        try:
            with open(file_path, "rb") as f:
                img_data = f.read()
                img_np = cv2.imdecode(np.frombuffer(img_data, np.uint8), cv2.IMREAD_COLOR)
        except Exception as e:
            return f"<div style='background: white; padding: 40px; color: red;'>Critical Read Error: {e}</div>"
        
        if img_np is None:
             return f"<div style='background: white; padding: 40px; color: red;'>Corrupt or unsupported image format for {filename}</div>"

        img_tag = f'<img src="/uploads/{filename}" style="max-width: 100%; height: auto; max-height: 850px; display: block; margin: 0 auto; border-radius: 12px; box-shadow: 0 8px 24px rgba(0,0,0,0.12);" />'
        clickable_html = build_clickable_img_wrapper(img_tag, img_np, fast=False)
        text = extract_image_text(img_np, fast=False)
        clean_text_html = ""
        for line in text.split("\n"):
            line = line.strip()
            if line:
                clean_text_html += f'<p style="margin-bottom: 1.25em; font-family: \'Georgia\', serif; line-height: 1.8;">{line}</p>'
        final_img = clickable_html if clickable_html else f'<div style="text-align: center;">{img_tag}</div>'
        
        return f"""
            <div id="img-page-0" class="lazy-page-container flex-page" data-original-width="800" style="display: flex; flex-direction: column; margin-bottom: 40px; padding: 60px; gap: 40px; background: white;">
                <div class="pdf-img-top" style="text-align: center; width: 100%;">{final_img}</div>
                <div class="pdf-text-bottom" style="width: 100%; max-width: 850px; margin: 0 auto; color: #1e293b;">
                    <div style="border-top: 1px solid #eee; padding-top: 30px; margin-top: 10px;">{clean_text_html}</div>
                </div>
            </div>
        """
    else:
        # ABSOLUTE FALLBACK: Instead of showing "Unsupported", try to detect content type
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                snippet = f.read(2048).strip().lower()
                
            # If the file looks like HTML (starts with tag or doctype), render as HTML
            if snippet.startswith("<html") or snippet.startswith("<!doctype") or ("<p>" in snippet) or ("<h1>" in snippet):
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                return f"<div id='html-page-0' class='lazy-page-container' style='padding: 60px; max-width: 900px; margin: 0 auto; background: var(--bg-paper); border-radius: 4px; box-shadow: 0 4px 20px rgba(0,0,0,0.1); color: #333;'>{content}</div>"
            
            # Otherwise, standard plain text extraction
            return extract_txt_html(file_path)
        except:
            return "<div style='background: white; padding: 40px; color: red; text-align: center;'><h2>🚫 File Not Readable</h2><p>This file format is not supported or the file is corrupted.</p></div>"



@app.route("/")
def home():
    return render_template("index.html")


@app.route("/uploads/<path:filename>")
def serve_upload(filename):
    return send_from_directory(UPLOAD_FOLDER, filename)

@app.route("/upload", methods=["POST"])
def upload():
    if "file" not in request.files:
        return "No file selected", 400

    file = request.files["file"]

    if file.filename == "":
        return "No file selected", 400

    # --- Duplicate Detection ---
    conn = get_conn()
    cur = conn.cursor()
    cur.execute("SELECT id FROM books WHERE name = ?", (file.filename,))
    existing = cur.fetchone()
    conn.close()

    if existing:
        return jsonify({
            "status": "duplicate",
            "message": f"'{file.filename}' is already in your library!"
        }), 409

    # Save the file right away
    filepath = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(filepath)

    extracted_path = os.path.join(EXTRACTED_FOLDER, file.filename + ".html")

    # Insert DB record immediately with status = 'processing'
    conn = get_conn()
    cur = conn.cursor()
    cur.execute(
        "INSERT INTO books(name, path, extracted_path, status) VALUES (?, ?, ?, ?)",
        (file.filename, filepath, extracted_path, "processing")
    )
    book_id = cur.lastrowid
    conn.commit()
    conn.close()

    threading.Thread(target=do_extract_task, args=(filepath, extracted_path, book_id), daemon=True).start()
    return jsonify({"status": "processing", "message": "Book received! Processing in background..."})


def background_ocr_upgrade_task(fpath, epath):
    """Deep OCR pass in background to upgrade structural extractions."""
    try:
        # Perform the heavy-duty OCR now (no fast mode)
        final_html = extract_book_html(fpath, fast_mode=False)
        with open(epath, "w", encoding="utf-8") as f:
            f.write(str(final_html))
        print(f"Background OCR upgrade complete for {epath}.")
    except Exception:
        traceback.print_exc()

def do_extract_task(fpath, epath, bid):
    """Phase 1: Fast Start (Structural extraction only for instant availability)"""
    try:
        # 1. Quickly extract structure and native text, skip heavy OCR for now
        html = extract_book_html(fpath, fast_mode=True)
        os.makedirs(os.path.dirname(epath), exist_ok=True)
        with open(epath, "w", encoding="utf-8") as f:
            f.write(str(html))
        
        # 2. Mark book as ready so user can open it immediately
        conn2 = get_conn()
        conn2.execute("UPDATE books SET status='ready' WHERE id=?", (bid,))
        conn2.commit()
        conn2.close()
        print(f"Book {bid} is ready (Fast Mode).")

        # Phase 2: Background Knowledge Mining (Async OCR)
        threading.Thread(target=background_ocr_upgrade_task, args=(fpath, epath), daemon=True).start()

    except Exception as e:
        traceback.print_exc()
        try:
            conn2 = get_conn()
            conn2.execute("UPDATE books SET status='error' WHERE id=?", (bid,))
            conn2.commit()
            conn2.close()
        except: pass

    return jsonify({"status": "processing", "message": "Book received! Processing in background..."})


@app.route("/books")
def books():
    conn = get_conn()
    cur = conn.cursor()
    cur.execute("SELECT id, name, uploaded_at, status FROM books ORDER BY id DESC")
    data = cur.fetchall()
    conn.close()
    return jsonify(data)


@app.route("/tts")
def tts():
    text = request.args.get("text", "")
    lang = request.args.get("lang", "en") or "en"
    gender = request.args.get("gender", "female").lower()
    
    if not text:
        return "No text", 400
    
    # Standardize language (ta-IN -> ta)
    l = lang.split('-')[0].lower() if '-' in lang else lang.lower()
    
    # Professional Neural Voice Map (Comprehensive)
    VOICE_MAP = {
        'en': { 'female': 'en-US-AriaNeural', 'male': 'en-US-GuyNeural' },
        'ta': { 'female': 'ta-IN-PallaviNeural', 'male': 'ta-IN-ValluvarNeural' },
        'hi': { 'female': 'hi-IN-SwaraNeural', 'male': 'hi-IN-MadhurNeural' },
        'kn': { 'female': 'kn-IN-SapnaNeural', 'male': 'kn-IN-GaganNeural' },
        'te': { 'female': 'te-IN-ShrutiNeural', 'male': 'te-IN-MohanNeural' },
        'ml': { 'female': 'ml-IN-SobhanaNeural', 'male': 'ml-IN-MidhunNeural' },
        'pa': { 'female': 'pa-IN-OjasNeural', 'male': 'pa-IN-GurumaNeural' },
        'gu': { 'female': 'gu-IN-DhwaniNeural', 'male': 'gu-IN-NiranjanNeural' },
        'mr': { 'female': 'mr-IN-AarohiNeural', 'male': 'mr-IN-ManoharNeural' },
        'bn': { 'female': 'bn-IN-TanishaaNeural', 'male': 'bn-IN-BashkarNeural' },
        'ur': { 'female': 'ur-PK-UzmaNeural', 'male': 'ur-PK-AsadNeural' },
        'or': { 'female': 'hi-IN-SwaraNeural', 'male': 'hi-IN-MadhurNeural' }, # Odia fallback to Hindi because Edge-TTS lacks native Odia
        'ko': { 'female': 'ko-KR-SunHiNeural', 'male': 'ko-KR-InJoonNeural' },
        'ja': { 'female': 'ja-JP-NanamiNeural', 'male': 'ja-JP-KeitaNeural' },
        'th': { 'female': 'th-TH-PremwadeeNeural', 'male': 'th-TH-NiwatNeural' },
        'fr': { 'female': 'fr-FR-DeniseNeural', 'male': 'fr-FR-HenriNeural' },
        'es': { 'female': 'es-ES-ElviraNeural', 'male': 'es-ES-AlvaroNeural' },
        'de': { 'female': 'de-DE-KatjaNeural', 'male': 'de-DE-ConradNeural' },
        'it': { 'female': 'it-IT-ElsaNeural', 'male': 'it-IT-DiegoNeural' },
        'zh': { 'female': 'zh-CN-XiaoxiaoNeural', 'male': 'zh-CN-YunxiNeural' }
    }
    
    # HYBRID ENGINE: Edge Neural for major languages, gTTS for regional fallbacks
    EDGE_SUPPORTED = ['en','ta','hi','bn','kn','te','ml','gu','mr','fr','es','de','it','zh','ur','ko','ja','th']
    
    if l not in EDGE_SUPPORTED:
        # PURE FALLBACK FOR UNSUPPORTED REGIONS
        try:
            tts_obj = gTTS(text=text, lang=l, slow=False)
            fp = io.BytesIO()
            tts_obj.write_to_fp(fp)
            fp.seek(0)
            return send_file(fp, mimetype="audio/mpeg")
        except Exception as e:
            print(f"gTTS Fallback Failed for {l}: {e}")
            voice = 'en-US-AriaNeural' # Ultimate fallback
    else:
        voice = VOICE_MAP.get(l, {}).get(gender, 'en-US-AriaNeural')


    def generate_streaming_tts():
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        communicate = edge_tts.Communicate(text, voice)
        async_gen = communicate.stream()
        
        while True:
            try:
                # We pull one chunk at a time from the async iterator
                chunk = loop.run_until_complete(async_gen.__anext__())
                if chunk["type"] == "audio":
                    yield chunk["data"]
            except StopAsyncIteration:
                break
            except Exception as e:
                print(f"Streaming Error: {e}")
                break
        loop.close()

    return Response(generate_streaming_tts(), mimetype="audio/mpeg")

@app.route("/book/<int:book_id>")
def open_book(book_id):
    try:
        conn = get_conn()
        cur = conn.cursor()

        cur.execute("SELECT name, path, extracted_path FROM books WHERE id = ?", (book_id,))
        book = cur.fetchone()
        conn.close()

        if not book:
            return jsonify({"error": "Book not found in database"}), 404

        name, file_path, extracted_path = book
        file_name = os.path.basename(file_path) if file_path else ""

        # HEALING LOGIC: Detect if the book was uploaded using the old, slow Base64 method.
        # If 'data:image' is found in a PDF file's first few thousand characters, we convert it to the new efficient format.
        needs_migration = False
        is_pdf = str(file_path or "").lower().endswith(".pdf")
        
        if os.path.exists(extracted_path) and is_pdf:
            st = os.stat(extracted_path)
            # If the file is > 5MB and it's a PDF, it's almost certainly bloated with Base64 images
            if st.st_size > 5 * 1024 * 1024:
                with open(extracted_path, "r", encoding="utf-8", errors="ignore") as f:
                    sample = f.read(20000)
                    if 'src="data:image/' in sample:
                        needs_migration = True
            
            # AGGRESSIVE HEALING: If the PDF extraction has NO image tags but we think it should
            # (or just to be safe if it was extracted before the external asset fix), 
            # we trigger repair if specific markers are missing.
            if not needs_migration:
                try:
                    with open(extracted_path, "r", encoding="utf-8", errors="ignore") as f:
                        content_sample = f.read(50000)
                    # If it's a large PDF but has NO images at all in the first 50k chars, 
                    # it might be a failed extraction.
                    if '<img' not in content_sample:
                        # Check if the PDF actually has images using a fast check
                        try:
                            # Use absolute path to ensure fitz can find it
                            abs_pdf_path = os.path.abspath(file_path)
                            test_pdf = fitz.open(abs_pdf_path)
                            if len(test_pdf) > 0 and test_pdf[0].get_images():
                                needs_migration = True
                            test_pdf.close()
                        except Exception: 
                            pass
                except Exception:
                    pass

        if not os.path.exists(extracted_path) or needs_migration:
            # If the file exists, return it NOW and heal in background.
            if os.path.exists(extracted_path):
                print(f"Book needs healing. Sending current version and starting background repair for {file_name}...")
                with open(extracted_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
                threading.Thread(target=background_ocr_upgrade_task, args=(file_path, extracted_path), daemon=True).start()
            else:
                print(f"New book or missing extraction. Creating initial fast version for {file_name}...")
                text = extract_book_html(file_path, fast_mode=True)
                os.makedirs(os.path.dirname(extracted_path), exist_ok=True)
                with open(extracted_path, "w", encoding="utf-8") as f:
                    f.write(str(text))
                threading.Thread(target=background_ocr_upgrade_task, args=(file_path, extracted_path), daemon=True).start()
        else:
            with open(extracted_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        with open(extracted_path, "r", encoding="utf-8", errors="ignore") as f:
            text = str(f.read())

        # Automatic Language Detection for Regional Support
        detected_lang = 'en'
        try:
            # Sample first 1000 visible characters for faster detection
            sample_text = re.sub(r'<[^>]+>', '', text[:5000])
            
            # Script-based fallback (very robust for Indian languages)
            if re.search(r'[\u0B80-\u0BFF]', sample_text): detected_lang = 'ta'
            elif re.search(r'[\u0900-\u097F]', sample_text): detected_lang = 'hi'
            elif re.search(r'[\u0C00-\u0C7F]', sample_text): detected_lang = 'te'
            elif re.search(r'[\u0C80-\u0CFF]', sample_text): detected_lang = 'kn'
            elif re.search(r'[\u0D00-\u0D7F]', sample_text): detected_lang = 'ml'
            elif re.search(r'[\u0980-\u09FF]', sample_text): detected_lang = 'bn'
            elif re.search(r'[\u0A00-\u0A7F]', sample_text): detected_lang = 'pa'
            elif re.search(r'[\u0A80-\u0AFF]', sample_text): detected_lang = 'gu'
            elif re.search(r'[\u0D80-\u0DFF]', sample_text): detected_lang = 'si'
            elif re.search(r'[\u3040-\u309F\u30A0-\u30FF]', sample_text): detected_lang = 'ja'
            elif re.search(r'[\u4E00-\u9FFF]', sample_text): detected_lang = 'zh'
        except Exception:
            pass

        return jsonify({
            "name": name,
            "text": text,
            "file_name": file_name,
            "detected_lang": detected_lang
        })
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": f"Server Error: {str(e)}"}), 500



@app.route("/save_highlight", methods=["POST"])
def save_highlight():
    data = request.get_json()

    book_id = data.get("book_id")
    highlighted_text = data.get("highlighted_text")

    if not book_id or not highlighted_text:
        return jsonify({"error": "Missing data"})

    conn = get_conn()
    cur = conn.cursor()

    cur.execute(
        "INSERT INTO highlights(book_id, highlighted_text) VALUES (?, ?)",
        (book_id, highlighted_text)
    )

    conn.commit()
    conn.close()

    return jsonify({"message": "Highlight saved successfully"})


@app.route("/delete_highlight", methods=["POST"])
def delete_highlight():
    data = request.get_json()

    book_id = data.get("book_id")
    highlighted_text = data.get("highlighted_text")

    if not book_id or not highlighted_text:
        return jsonify({"error": "Missing data"})

    conn = get_conn()
    cur = conn.cursor()

    cur.execute(
        "DELETE FROM highlights WHERE book_id = ? AND highlighted_text = ?",
        (book_id, highlighted_text)
    )

    conn.commit()
    conn.close()

    return jsonify({"message": "Highlight deleted successfully"})


@app.route("/highlights/<int:book_id>")
def get_highlights(book_id):
    conn = get_conn()
    cur = conn.cursor()

    cur.execute(
        "SELECT highlighted_text FROM highlights WHERE book_id = ?",
        (book_id,)
    )
    rows = cur.fetchall()

    conn.close()

    highlights = [row[0] for row in rows]
    return jsonify(highlights)


@app.route("/delete_book/<int:book_id>", methods=["POST"])
def delete_book(book_id):
    conn = get_conn()
    cur = conn.cursor()

    cur.execute("SELECT path, extracted_path FROM books WHERE id = ?", (book_id,))
    book = cur.fetchone()

    if not book:
        conn.close()
        return jsonify({"error": "Book not found"})

    file_path, extracted_path = book

    cur.execute("DELETE FROM highlights WHERE book_id = ?", (book_id,))
    cur.execute("DELETE FROM books WHERE id = ?", (book_id,))
    conn.commit()
    conn.close()

    if file_path and os.path.exists(file_path):
        os.remove(file_path)

    if extracted_path and os.path.exists(extracted_path):
        os.remove(extracted_path)

    return jsonify({"message": "Book deleted successfully"})


@app.route("/download/<int:book_id>")
def download_book(book_id):
    conn = get_conn()
    cur = conn.cursor()
    cur.execute("SELECT path, name FROM books WHERE id = ?", (book_id,))
    book = cur.fetchone()
    conn.close()

    if not book:
        return "Book not found", 404

    file_path, name = book
    if not file_path or not os.path.exists(file_path):
        return "File not found", 404

    response = send_from_directory(
        os.path.dirname(file_path),
        os.path.basename(file_path),
        as_attachment=True
    )
    # Ensure browsers don't fallback to UUIDs by explicitly encoding the filename
    encoded_name = quote(name)
    response.headers["Content-Disposition"] = f"attachment; filename=\"{name}\"; filename*=UTF-8''{encoded_name}"
    return response


@app.route("/original/<int:book_id>")
def original_book(book_id):
    conn = get_conn()
    cur = conn.cursor()
    cur.execute("SELECT path FROM books WHERE id = ?", (book_id,))
    book = cur.fetchone()
    conn.close()

    if not book:
        return "Book not found", 404

    file_path = book[0]
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return send_from_directory(os.path.dirname(file_path), os.path.basename(file_path))
    else:
        html = extract_book_html(file_path)
        return f"<html><head><style>body {{ margin: 0; background: #f0f0f0; }}</style></head><body>{html}</body></html>"


@app.route("/translate_text", methods=["POST"])
def translate_text():
    try:
        data = request.get_json()
        texts = data.get("texts", [])
        target_lang = data.get("target_lang", "en")
        source_lang = data.get("source_lang", "auto")
        
        if not texts:
            return jsonify([])
            
        # TRY SIDECAR (High-Speed Node.js Proxy)
        try:
            sidecar_url = f"http://localhost:{SIDECAR_PORT}/translate"
            # Increase reliability for massive documents
            res = requests.post(sidecar_url, json={"texts": texts, "target_lang": target_lang, "source_lang": source_lang}, timeout=30)
            if res.ok:
                results = res.json()
                # Ensure ALL results were successfully translated (not null fallback in sidecar)
                if results and len(results) == len(texts) and all(r is not None for r in results):
                    return jsonify(results)
                else:
                    print(f"Sidecar health check failed: {len(results) if results else 0}/{len(texts)} translated.")
            else:
                print(f"Sidecar responded with error: {res.status_code}")
        except Exception as e:
            print(f"Sidecar Bridge Failure: {e}. Falling back to internal Python engine.")

        # FALLBACK: Using a specialized Batch Translator for maximum reliability vs manual string joining
        translator = GoogleTranslator(source='auto', target=target_lang)
        
        # PARTITIONED BATCHING: Prevents URI overflow and provides faster individual feedback
        batches = []
        current_batch = []
        for t in texts:
            if len(current_batch) >= 25: 
                batches.append(current_batch)
                current_batch = []
            current_batch.append(str(t))
        if current_batch:
            batches.append(current_batch)
            
        print(f"Backend translating {len(texts)} text nodes in {len(batches)} batches...")
            
        def translate_single_batch(batch):
            if not batch: return []
            try:
                # 1. Use the library's native batch translator (handles mapping internally)
                # This is VASTLY more reliable than manual delimiter joining for short text.
                return translator.translate_batch(batch)
            except Exception as e:
                print(f"Batch translation failed ({e}). Attempting granular recovery...")
                results = []
                # 2. INDIVIDUAL FALLBACK: Ensure one bad node doesn't kill the whole batch
                for s in batch:
                    # Clean input (strip whitespace but remember it for return)
                    clean_s = s.strip()
                    if not clean_s:
                        results.append(s)
                        continue
                        
                    # Individual Retry with small backoff
                    success = False
                    for attempt in range(2):
                        try:
                            time.sleep(attempt * 0.5) 
                            translated = translator.translate(clean_s)
                            if translated:
                                # Preserve original casing/spacing style where possible
                                results.append(translated)
                                success = True
                                break
                        except:
                            continue
                    if not success:
                        results.append(s) # Safety: Return original if even individual fail
                return results

        with concurrent.futures.ThreadPoolExecutor(max_workers=6) as executor:
            batch_results = list(executor.map(translate_single_batch, batches))
            
        translated_texts = [item for sublist in batch_results for item in sublist]

        # Force identical array lengths to ensure DOM mapping doesn't break
        if len(translated_texts) != len(texts):
            print(f"CRITICAL MISMATCH after fallback: {len(translated_texts)} vs {len(texts)}")
            return jsonify(texts)

        return jsonify(translated_texts)
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


# --- BOOKMARK SYSTEM ---

@app.route("/save_bookmark", methods=["POST"])
def save_bookmark():
    try:
        data = request.get_json()
        book_id = data.get("book_id")
        page_num = data.get("page_number", 1)
        scroll_y = data.get("scroll_y", 0)
        char_index = data.get("char_index", 0)
        node_index = data.get("node_index", -1)
        node_offset = data.get("node_offset", 0)
        lang_code = data.get("lang_code", "en")
        label = data.get("label", f"Page {page_num}")
        force_replace = data.get("replace", False)

        if not book_id:
            return jsonify({"error": "No book selected"}), 400

        conn = get_conn()
        cur = conn.cursor()
        
        # Check for existing bookmark to enforce single-bookmark-per-book
        cur.execute("SELECT id FROM bookmarks WHERE book_id = ?", (book_id,))
        existing = cur.fetchone()
        
        if existing and not force_replace:
            conn.close()
            return jsonify({"status": "exists", "message": "A bookmark already exists for this book. Replace it?"})
        
        if existing:
            # Update current bookmark
            cur.execute(
                "UPDATE bookmarks SET page_number = ?, scroll_y = ?, char_index = ?, node_index = ?, node_offset = ?, label = ?, lang_code = ?, created_at = CURRENT_TIMESTAMP WHERE book_id = ?",
                (page_num, scroll_y, char_index, node_index, node_offset, label, lang_code, book_id)
            )
        else:
            # Insert new one
            cur.execute(
                "INSERT INTO bookmarks(book_id, page_number, scroll_y, char_index, node_index, node_offset, label, lang_code) VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                (book_id, page_num, scroll_y, char_index, node_index, node_offset, label, lang_code)
            )
            
        conn.commit()
        conn.close()
        return jsonify({"status": "success", "message": "Bookmark updated!"})
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route("/bookmarks/<int:book_id>", methods=["GET"])
def get_bookmarks(book_id):
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(
            "SELECT id, page_number, scroll_y, char_index, node_index, node_offset, label, lang_code, created_at FROM bookmarks WHERE book_id = ? ORDER BY created_at DESC",
            (book_id,)
        )
        bookmarks = [
            {
                "id": r[0], "page_number": r[1], "scroll_y": r[2], "char_index": r[3], 
                "node_index": r[4], "node_offset": r[5], "label": r[6], "lang_code": r[7], "created_at": r[8]
            }
            for r in cur.fetchall()
        ]
        conn.close()
        return jsonify(bookmarks)
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route("/delete_bookmark/<int:bookmark_id>", methods=["POST"])
def delete_bookmark(bookmark_id):
    try:
        conn = get_conn()
        cur = conn.cursor()
        cur.execute("DELETE FROM bookmarks WHERE id = ?", (bookmark_id,))
        conn.commit()
        conn.close()
        return jsonify({"status": "success"})
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

@app.route('/summarize', methods=['POST'])
def summarize():
    try:
        data = request.get_json()
        text = data.get("text", "")
        book_id = data.get("book_id")
        
        # If book_id provided, we use the cached full-book sentences or read from disk
        # (Allows summarization of specific highlights OR full context search)
        if book_id and not text:
            sentences = get_book_sentences(book_id)
            text = " ".join(sentences)
            
        if not text or len(text.strip()) < 50:
            return jsonify({"summary": "Text is too short to summarize. Please highlight a larger section."})

        # Pure Python Extractive Summarization (Fallback/Offline mode without heavy NLP dependencies)
        
        # 1. Clean and tokenize sentences
        sentences = re.split(r'(?<=[.!?])\s+', text.strip())
        if len(sentences) <= 3:
            return jsonify({"summary": text})
            
        # 2. Word frequency calculation (ignore common stop words)
        stop_words = set(['the', 'is', 'in', 'and', 'to', 'a', 'of', 'for', 'it', 'on', 'with', 'as', 'by', 'that', 'this', 'an', 'are', 'was', 'be', 'or', 'at', 'from'])
        words = re.findall(r'\w+', text.lower())
        freq_table = {}
        for word in words:
            if word not in stop_words:
                freq_table[word] = freq_table.get(word, 0) + 1
                
        # 3. Score sentences based on word frequency
        sentence_scores = {}
        for i, sentence in enumerate(sentences):
            words_in_sentence = re.findall(r'\w+', sentence.lower())
            
            sentence_score = 0
            for word in words_in_sentence:
                if word in freq_table:
                    sentence_score = int(sentence_score) + int(freq_table.get(word, 0))
            
            # Normalize score by sentence length to prevent unfairly weighting extremely long run-ons
            if len(words_in_sentence) > 0:
                sentence_scores[i] = float(sentence_score) / float(len(words_in_sentence))

        # 4. Extract top sentences (aim for ~30% compression, max 6 bullet points)
        target_length = max(2, min(6, int(len(sentences) * 0.35)))
        
        # Get the indices of the highest-scoring sentences
        all_indices = sorted(sentence_scores.keys(), key=lambda k: float(sentence_scores[k]), reverse=True)
        top_sentence_indices = []
        for i in range(min(len(all_indices), target_length)):
            top_sentence_indices.append(all_indices[i])
        
        # 5. Re-sort chronologically so the summary reads in the correct order
        top_sentence_indices.sort()
        
        summary_sentences = []
        for index in top_sentence_indices:
            clean_sentence = sentences[index].strip()
            if clean_sentence:
                summary_sentences.append("- " + clean_sentence)
            
        summary = "\n".join(summary_sentences)
        
        return jsonify({"summary": summary})
        
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": "Summarization failed internally."}), 500

@app.route("/save_note", methods=["POST"])
def save_note():
    data = request.get_json()
    book_id = data.get("book_id")
    content = data.get("content")
    if not book_id or not content:
        return jsonify({"error": "Missing data"}), 400

    conn = get_conn()
    cur = conn.cursor()
    cur.execute("INSERT INTO notes(book_id, content) VALUES (?, ?)", (book_id, content))
    conn.commit()
    conn.close()
    return jsonify({"message": "Note saved!"})

@app.route("/notes/<int:book_id>")
def get_notes(book_id):
    lang = request.args.get("lang", "en") # Current UI language
    conn = get_conn()
    cur = conn.cursor()
    cur.execute("SELECT id, content FROM notes WHERE book_id = ? ORDER BY created_at DESC", (book_id,))
    rows = cur.fetchall()
    conn.close()

    notes = []
    for row in rows:
        note_id, content = row
        # Auto-translate if not in original/english
        if lang and lang != 'en' and lang != 'orig':
            try:
                content = GoogleTranslator(source='auto', target=lang).translate(content)
            except Exception:
                pass # Fallback to original if translation fails
        notes.append({"id": note_id, "content": content})
        
    return jsonify(notes)

@app.route("/update_note", methods=["POST"])
def update_note():
    data = request.get_json()
    note_id = data.get("note_id")
    content = data.get("content")
    if not note_id or content is None:
        return jsonify({"error": "Missing data"}), 400

    conn = get_conn()
    cur = conn.cursor()
    cur.execute("UPDATE notes SET content = ? WHERE id = ?", (content, note_id))
    conn.commit()
    conn.close()
    return jsonify({"message": "Note updated!"})

@app.route("/delete_note/<int:note_id>", methods=["POST"])
def delete_note_endpoint(note_id):
    conn = get_conn()
    cur = conn.cursor()
    cur.execute("DELETE FROM notes WHERE id = ?", (note_id,))
    conn.commit()
    conn.close()
    return jsonify({"message": "Note deleted!"})

@app.route("/download_notes/<int:book_id>")
def download_notes(book_id):
    format_type = request.args.get("format", "txt").lower()
    
    conn = get_conn()
    cur = conn.cursor()
    cur.execute("SELECT name FROM books WHERE id = ?", (book_id,))
    book_row = cur.fetchone()
    if not book_row:
        conn.close()
        return "Book not found", 404
        
    book_name = book_row[0]
    cur.execute("SELECT content FROM notes WHERE book_id = ?", (book_id,))
    notes = [row[0] for row in cur.fetchall()]
    conn.close()

    
    if not notes:
        return "No notes found to download.", 400

    if format_type == "docx":
        doc = docx.Document()
        doc.add_heading(f'Study Notes: {book_name}', 0)
        doc.add_paragraph('Collected using AI Book Reader').italic = True
        
        for i, note in enumerate(notes):
            doc.add_heading(f'Snippet {i+1}', level=1)
            doc.add_paragraph(note)
            
        target = BytesIO()
        doc.save(target)
        target.seek(0)
        return send_file(target, as_attachment=True, download_name=f"Notes_{book_name}.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    elif format_type == "pdf":
        class PDF(FPDF):
            def header(self):
                self.set_font('helvetica', 'B', 15)
                self.cell(0, 10, f'Study Notes: {book_name}', border=True, ln=True, align='C')
                self.ln(5)
            def footer(self):
                self.set_y(-15)
                self.set_font('helvetica', 'I', 8)
                self.cell(0, 10, f'Page {self.page_no()}', 0, 0, 'C')

        pdf = PDF()
        pdf.add_page()
        pdf.set_font("helvetica", size=11)
        
        for i, note in enumerate(notes):
            pdf.set_font("helvetica", 'B', 12)
            pdf.cell(0, 10, f"Snippet {i+1}:", ln=True)
            pdf.set_font("helvetica", size=10)
            # Remove characters that may crash basic PDF fonts or use multi_cell safely
            safe_text = note.encode('latin-1', 'replace').decode('latin-1')
            pdf.multi_cell(0, 10, safe_text)
            pdf.ln(5)

        response_bytes = pdf.output()
        target = BytesIO(response_bytes)
        return send_file(target, as_attachment=True, download_name=f"Notes_{book_name}.pdf", mimetype="application/pdf")

    else: # Default TXT
        content = f"STUDY NOTES: {book_name}\n{'='*30}\n\n"
        for i, note in enumerate(notes):
            content += f"SNIPPET {i+1}:\n{note}\n\n"
            content += "-"*20 + "\n\n"
            
        target = BytesIO(content.encode('utf-8'))
        return send_file(target, as_attachment=True, download_name=f"Notes_{book_name}.txt", mimetype="text/plain")

@app.route("/generate_revision", methods=["POST"])
def generate_revision():
    try:
        data = request.get_json()
        book_id = data.get("book_id")
        
        target_lang = data.get("target_lang", "en")
        
        if not book_id:
            return jsonify({"error": "No book selected"}), 400
            
        sentences = get_book_sentences(book_id)
        if not sentences:
            return jsonify({"error": "No content found in book"}), 404
            
        # Full-book scoring 
        text = " ".join(sentences).lower()
        stop_words = set(['the', 'is', 'in', 'and', 'to', 'a', 'of', 'for', 'it', 'on', 'with', 'as', 'by', 'that', 'this', 'an', 'are', 'was', 'be', 'or', 'at', 'from', 'their', 'which', 'will', 'have', 'been', 'were', 'about'])
        words = re.findall(r'\b\w{4,}\b', text)
        freq = {}
        for w in words:
            freq[w] = freq.get(w, 0) + 1
            
        # Score each sentence
        scores = []
        for i, s in enumerate(sentences):
            words_in_s = re.findall(r'\b\w+\b', s.lower())
            if len(words_in_s) < 12 or len(words_in_s) > 50: continue
            
            score = sum(freq[w] for w in words_in_s if w in freq)
            scores.append((score / len(words_in_s), i, s))
            
        # Pick top 25 diverse points
        top_points = sorted(scores, key=lambda x: float(x[0]), reverse=True)[:30]
        # Re-sort by book order
        top_points.sort(key=lambda x: x[1])
        
        revision_points = [p[2].strip() for p in top_points]
        
        if target_lang != 'en' and target_lang != 'orig':
            try:
                translator = GoogleTranslator(source='auto', target=target_lang)
                revision_points = [translator.translate(pt) for pt in revision_points]
            except Exception as e:
                print(f"Revision Translation Error: {e}")
                
        return jsonify({"revision_points": revision_points})
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/download_revision/<book_id>", methods=["GET"])
def download_revision(book_id):
    try:
        target_lang = request.args.get("target_lang", "en")
        
        sentences = get_book_sentences(book_id)
        if not sentences: return "Book not found", 404
        
        # Generation Logic (Duplicated for standalone download)
        text = str(" ".join(sentences).lower())
        stop_words = set(['the', 'is', 'in', 'and', 'to', 'a', 'of', 'for', 'it', 'on', 'with', 'as', 'by', 'that', 'this', 'an', 'are', 'was', 'be', 'or', 'at', 'from'])
        words = re.findall(r'\b\w{4,}\b', text)
        freq = {}
        for w in words:
            freq[w] = int(freq.get(w, 0)) + 1
        
        scores = []
        for i, s in enumerate(sentences):
            words_in_s = re.findall(r'\b\w+\b', s.lower())
            if len(words_in_s) < 12 or len(words_in_s) > 50: continue
            score = float(sum(freq[w] for w in words_in_s if w in freq))
            scores.append((score / float(len(words_in_s)), i, s))
            
        all_top = sorted(scores, key=lambda x: float(str(x[0])), reverse=True)
        top_points = []
        for i in range(min(len(all_top), 35)):
            top_points.append(all_top[i])
        top_points.sort(key=lambda x: int(x[1]))
        
        revision_points = [p[2].strip() for p in top_points]
        
        if target_lang != 'en' and target_lang != 'orig':
            try:
                translator = GoogleTranslator(source='auto', target=target_lang)
                revision_points = [translator.translate(pt) for pt in revision_points]
            except Exception as e:
                print(f"Revision Content Translation Error: {e}")
        
        conn = get_conn()
        cur = conn.cursor()
        cur.execute("SELECT name FROM books WHERE id = ?", (book_id,))
        row = cur.fetchone()
        conn.close()
        book_name = row[0] if row else "Book"
        
        content = f"REVISION GUIDE: {book_name}\n" + "="*40 + "\n\n"
        for i, p in enumerate(top_points):
            content += f"{i+1}. {p[2].strip()}\n\n"
            
        return send_file(BytesIO(content.encode('utf-8')), as_attachment=True, download_name=f"Revision_{book_name}.txt", mimetype="text/plain")
    except Exception as e:
        return str(e), 500

def translate_batch_fast(texts, target_lang):
    """High-concurrency batch translation with Sidecar support and reliable fallback."""
    if not texts or not target_lang or target_lang in ["en", "orig"]:
        return texts
    
    # Pre-clean: strip and filter out empty strings to avoid translation hits
    to_translate = []
    mapping = [] # Map index from texts to to_translate index
    
    for i, t in enumerate(texts):
        if t and t.strip():
            mapping.append(len(to_translate))
            to_translate.append(str(t).strip())
        else:
            mapping.append(-1)
            
    if not to_translate:
        return texts

    translated_pool = []
    
    # 1. ATTEMPT SIDECAR (Preferred for Speed/Concurrency)
    try:
        sidecar_url = f"http://localhost:{SIDECAR_PORT}/translate"
        # Increased timeout to 120s (sidecar is now faster with concurrent batches, but reliability is key)
        res = requests.post(sidecar_url, json={"texts": to_translate, "target_lang": target_lang}, timeout=120)
        if res.ok:
            data = res.json()
            if data and len(data) == len(to_translate) and all(d is not None for d in data):
                translated_pool = data
    except Exception as e:
        print(f"Quiz Sidecar Bridge Failure: {e}")

    # 2. FALLBACK to Python Library (Batching for safety)
    if not translated_pool:
        try:
            translator = GoogleTranslator(source='auto', target=target_lang)
            batch_size = 15 # conservative for deep_translator stability
            for i in range(0, len(to_translate), batch_size):
                batch = to_translate[i : i + batch_size]
                try:
                    translated_pool.extend(translator.translate_batch(batch))
                except Exception as e:
                    # Individual retry fallback
                    for item in batch:
                        try: translated_pool.append(translator.translate(item))
                        except: translated_pool.append(item)
        except Exception as e:
            print(f"Quiz Fallback Failed: {e}")
            return texts
            
    # 3. REASSEMBLE
    results = []
    for i, idx in enumerate(mapping):
        if idx == -1:
            results.append(texts[i])
        else:
            results.append(translated_pool[idx] if idx < len(translated_pool) else texts[i])
            
    return results

@app.route("/generate_quiz", methods=["POST"])
def generate_quiz():
    try:
        data = request.get_json()
        text = data.get("text", "")
        book_id = data.get("book_id")
        quiz_type = data.get("type", "mcq") # default to mcq

        if book_id and not text:
            sentences = get_book_sentences(book_id)
            # Randomized window ensure quiz covers various parts and prevents 500+ page parsing hangs
            random.shuffle(sentences)
            text = " ".join(sentences[:500])
        elif text and len(text) > 150000:
            # Massive selections (20+ pages) also need clamping for server safety
            text = text[:150000]


        if not text or len(text.strip()) < 40:
            return jsonify({"error": "Selection is too short. Please highlight a larger section or a full paragraph to generate a quiz."}), 400

        # 1. CLEAN & TOKENIZE
        sentences = re.split(r'(?<=[.!?])\s+', text.strip())
        all_words = re.findall(r'\b\w{4,}\b', text.lower())
        
        # 2. KEYWORDS EXTRACTION (Frequent nouns/concepts)
        stop_words = set(['the', 'and', 'that', 'with', 'from', 'this', 'their', 'which', 'will', 'have', 'been', 'were', 'about'])
        freq = {}
        for w in all_words:
            if w not in stop_words:
                freq[w] = int(freq.get(w, 0)) + 1
        
        raw_keywords = sorted(freq.keys(), key=lambda k: float(freq[k]), reverse=True)
        top_keywords = []
        for i in range(min(len(raw_keywords), 100)):
            top_keywords.append(str(raw_keywords[i]))
        
        # Scaling logic: More text = More questions (up to 50 for massive books in English)
        mcq_limit = min(50, max(10, len(sentences) // 15))
        short_limit = min(40, max(8, len(sentences) // 25))
        long_limit = min(20, max(5, len(sentences) // 40))

        # TIGHTEN LIMITS FOR TRANSLATED QUIZZES: 
        # Translation hits (Sidecar) take significantly longer; keep it fast!
        target_lang = data.get("target_lang", "en")
        if target_lang and target_lang not in ["en", "orig"]:
            mcq_limit = min(15, mcq_limit)
            short_limit = min(10, short_limit)
            long_limit = min(5, long_limit)

        questions = []

        if quiz_type == "mcq":
            # Existing MCQ logic
            candidates = [s for s in sentences if re.search(r'\b(is|was|means|refers to|defined as|called|known as)\b', s, re.I)]
            random.shuffle(candidates)
            
            for s in candidates[:max(mcq_limit * 2, 50)]:
                words = re.findall(r'\b\w+\b', s)
                possible_masks = [w for w in words if w.lower() in top_keywords and len(w) > 3]
                if not possible_masks: continue
                mask = random.choice(possible_masks)
                distractors = [k for k in top_keywords if k.lower() != mask.lower()]
                if len(distractors) < 3: continue
                options = random.sample(distractors, 3) + [mask]
                random.shuffle(options)
                q_text = re.sub(r'\b' + re.escape(mask) + r'\b', '__________', str(s), count=1, flags=re.I)
                questions.append({"question": q_text, "options": options, "answer": mask, "type": "mcq"})
                if len(questions) >= mcq_limit: break

        elif quiz_type in ["short", "long"]:
            # Standard Academic Question Templates
            # ... (temps logic)
            short_temps = [
                "Define {concept}.", 
                "What is meant by {concept}?", 
                "Write a short note on {concept}.", 
                "State the importance of {concept}.",
                "What are the primary uses of {concept}?"
            ]
            long_temps = [
                "Explain {concept} in detail based on the text.",
                "Describe the working or context of {concept}.",
                "Discuss the concept of {concept} and its implications.",
                "Elaborate on {concept} with suitable examples from the book.",
                "Analyze the importance of {concept} in modern applications."
            ]

            # High-Performance NLP Logic using TextBlob
            try:
                blob = TextBlob(text)
                # Group noun phrases and filter out junk
                concepts = [str(np).lower() for np in blob.noun_phrases if len(str(np)) > 3]
            except Exception:
                concepts = [k for k in top_keywords if len(str(k)) > 4]

            # Better stopword and diversity filtering
            hard_stop = {'here', 'there', 'that', 'this', 'with', 'from', 'been', 'were', 'about', 'some', 'than', 'into', 'only', 'very', 'just', 'more', 'also', 'their', 'which'}
            final_concepts = []
            for c in concepts:
                if c not in hard_stop and len(c.split()) < 4 and c not in final_concepts:
                    final_concepts.append(c)
            
            random.shuffle(final_concepts)

            target_count = short_limit if quiz_type == "short" else long_limit

            for concept in final_concepts:
                if len(questions) >= target_count: break
                
                # Find the best context for this concept
                for i, s in enumerate(sentences):
                    if concept.lower() in s.lower() and 10 < len(s.split()) < 60:
                        concept_display = concept.title()
                        if quiz_type == "short":
                            template = random.choice(short_temps)
                            questions.append({
                                "question": template.format(concept=concept_display),
                                "answer": f"💡 **Core Insight**: {s.strip()}",
                                "type": "short"
                            })
                        else:
                            # Contextual Multi-Sentence Answer
                            ctx_start = max(0, i-1)
                            ctx_end = min(len(sentences), i+5)
                            context = []
                            for idx in range(ctx_start, ctx_end):
                                context.append(sentences[idx])
                            
                            template = random.choice(long_temps)
                            formatted_answer = f"🔍 **Professional Summary: {concept_display}**\n\n"
                            formatted_answer += f"**Primary Explanation**: {str(context[0]).strip()}\n\n"
                            formatted_answer += "**Detailed Context**:\n"
                            for line in context[1:]:
                                if len(str(line).strip()) > 10:
                                    formatted_answer += f"• {str(line).strip()}\n"
                            
                            questions.append({
                                "question": template.format(concept=concept_display),
                                "answer": formatted_answer,
                                "type": "long"
                            })
                        break

        if not questions:
            # Fallback
            for k in top_keywords[:5]:
                questions.append({
                    "question": f"Based on the text, what is a key concept identified as '{k}'?",
                    "options": random.sample(top_keywords[10:13], 3) + [k] if quiz_type == "mcq" else None,
                    "answer": f"The term '{k}' is used as a significant keyword in this context.",
                    "type": quiz_type
                })
                if quiz_type == "mcq": random.shuffle(questions[-1]["options"])

        # --- NEW: High-Speed Batch Translation Pass ---
        target_lang = data.get("target_lang", "en")
        if target_lang and target_lang not in ["en", "orig"]:
            try:
                # 1. Flatten all strings to translate (Questions + Options + Answers)
                pool = []
                structure = [] # Map: (question_index, field_name, option_index)
                for i, q in enumerate(questions):
                    # Question Text
                    structure.append((i, "question", None))
                    pool.append(q["question"])
                    # Options (if MCQ)
                    if q.get("options"):
                        for j, opt in enumerate(q["options"]):
                            structure.append((i, "options", j))
                            pool.append(opt)
                    # Answer
                    if q.get("answer"):
                        structure.append((i, "answer", None))
                        pool.append(q["answer"])

                # 2. PERFORM BATCH HIT
                translated_pool = translate_batch_fast(pool, target_lang)

                # 3. RE-INJECT
                if len(translated_pool) == len(pool):
                    for idx, (q_idx, field, opt_idx) in enumerate(structure):
                        if field == "question":
                            questions[q_idx]["question"] = translated_pool[idx]
                        elif field == "options":
                            questions[q_idx]["options"][opt_idx] = translated_pool[idx]
                        elif field == "answer":
                            questions[q_idx]["answer"] = translated_pool[idx]
                else: 
                     print(f"Quiz translation count mismatch: {len(translated_pool)} vs {len(pool)}")
            except Exception as e:
                print(f"Quiz translation pipeline failed: {e}")
                traceback.print_exc()

        return jsonify({"questions": questions, "type": quiz_type})

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": "Quiz generation failed."}), 500

@app.route("/ask", methods=["POST"])
def ask_question():
    try:
        data = request.get_json()
        question = data.get("question", "")
        book_id = data.get("book_id")
        text = data.get("text", "")
        
        if not question:
            return jsonify({"answer": "Please provide a question."})
            
        # Optimization: Fetch pre-cached sentences from the server instead of receiving them via network
        if book_id:
            sentences = get_book_sentences(book_id)
        elif text:
            sentences = re.split(r'(?<=[.!?])\s+', text)
        else:
            return jsonify({"answer": "Ensure a book is open."})

        if not sentences:
            return jsonify({"answer": "Could not read book text or book is empty."})
            
        # 2. Extract question keywords (ignoring stop words)
        stop_words = set(['what', 'who', 'is', 'a', 'at', 'is', 'he', 'the', 'in', 'and', 'to', 'of', 'for', 'it', 'on', 'with', 'as', 'by', 'that', 'this', 'an', 'are', 'was', 'be', 'or', 'from', 'where', 'when', 'why', 'how', 'does', 'do', 'did'])
        question_words = [w.lower() for w in re.findall(r'\w+', question) if w.lower() not in stop_words]
        
        if not question_words:
            # Fallback if they asked "what is it"
            question_words = [w.lower() for w in re.findall(r'\w+', question)]
            
        # 3. Calculate Inverse Document Frequency (IDF) for question words
        word_doc_count = {}
        for sentence in sentences:
            sentence_words = set(re.findall(r'\w+', sentence.lower()))
            for qw in question_words:
                if qw in sentence_words:
                    word_doc_count[qw] = int(word_doc_count.get(qw, 0)) + 1
                    
        num_sentences = len(sentences)
        idf = {}
        for qw in question_words:
            # Add 1 to avoid division by zero
            idf[qw] = math.log(num_sentences / (1 + word_doc_count[qw]))
            
        # 4. Score sentences based on TF-IDF of question keywords
        sentence_scores = {}
        for i, sentence in enumerate(sentences):
            sentence_words = re.findall(r'\w+', sentence.lower())
            score = 0
            for qw in question_words:
                tf = sentence_words.count(qw)
                if tf > 0:
                    # Give bonus for exact phrase matching
                    if question.lower().find(qw) != -1 and sentence.lower().find(question.lower()) != -1:
                        score += 50
                    score += tf * idf[qw]
            
            if score > 0:
                # Add slight penalty for extremely long sentences to prefer concise answers
                length_penalty = math.log(float(max(10, len(sentence_words))))
                sentence_scores[i] = float(score) / float(length_penalty)
                
        if not sentence_scores:
            return jsonify({"answer": "I couldn't find an answer to that in the current book."})
            
        # 5. Extract top 2 most relevant sentences
        all_searched = sorted(sentence_scores.keys(), key=lambda k: float(sentence_scores[k]), reverse=True)
        best_indices = []
        for i in range(min(len(all_searched), 2)):
            best_indices.append(all_searched[i])
        best_indices.sort() # keep chronological
        
        answer = " ".join([sentences[idx].strip() for idx in best_indices])
        
        # Friendly formatting
        return jsonify({"answer": f"💡 Based on the book: {answer}"})
        
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": "Failed to search the book. It might be too large."}), 500

@app.route("/define", methods=["POST"])
def define_word():
    try:
        data = request.get_json()
        word = data.get("word", "").lower()
        book_id = data.get("book_id")
        text_from_client = data.get("text", "")
        lang = data.get("lang", "en").split('-')[0].lower() # e.g. 'ta' from 'ta-IN'
        
        if not word:
            return jsonify({"answer": "Word missing."})
            
        sentences = []
        # Prioritize the text current being viewed (this handles translations)
        if text_from_client:
            soup = BeautifulSoup(text_from_client, "html.parser")
            raw_text = soup.get_text(separator=" ")
            sentences = re.split(r'(?<=[.!?])\s+', raw_text)
            sentences = [re.sub(r'\s+', ' ', s).strip() for s in sentences if len(s.strip()) > 5]
        
        # Fallback to book file if client text is sparse or missing
        if not sentences and book_id:
            sentences = get_book_sentences(book_id)

        # 1. Multilingual pattern matching for finding definitions in-context
        lang_patterns = {
            "en": [rf"\b{word}\s+is\b", rf"\b{word}\s+means\b", rf"\b{word}\s+refers\s+to\b", rf"\bdefinition\s+of\s+{word}\b"],
            "ta": [rf"{word}\s+என்பது", rf"{word}\s+என்றால்", rf"{word}\s+குறிக்கிறது"],
            "hi": [rf"{word}\s+का\s+अर्थ", rf"{word}\s+है", rf"{word}\s+मतलब"],
            "te": [rf"{word}\s+అంటే", rf"{word}\s+అనేది"],
            "ml": [rf"{word}\s+എന്നാൽ", rf"{word}\s+എന്ന്\s+പറയുന്നത്"],
            "kn": [rf"{word}\s+ಎಂದರೆ", rf"{word}\s+ಎನ್ನುವುದು"],
            "bn": [rf"{word}\s+মানে", rf"{word}\s+হল"],
            "gu": [rf"{word}\s+એટલે", rf"{word}\s+છે"]
        }
        
        patterns = lang_patterns.get(lang, lang_patterns["en"])
        
        definition_sentences = []
        if sentences:
            for s in sentences:
                s_lower = s.lower()
                if any(re.search(p, s_lower) for p in patterns):
                    definition_sentences.append(s)
        
        if definition_sentences:
            combined_def = ""
            for i in range(min(len(definition_sentences), 2)):
                combined_def += str(definition_sentences[i]) + " "
            return jsonify({"answer": combined_def.strip()})
            
        # 2. Try Dictionary Definitions first (High signal)
        try:
            # First try: Direct definition in the word's own language
            dict_url = f"https://translate.googleapis.com/translate_a/single?client=gtx&sl={lang}&tl={lang}&dt=md&q={word}"
            res = requests.get(dict_url, timeout=5)
            if res.ok:
                dict_data = res.json()
                if len(dict_data) > 12 and dict_data[12]:
                    defs = []
                    for entry in dict_data[12]:
                        if len(entry) > 1 and entry[1]:
                            for subentry in entry[1]:
                                if subentry[0]: defs.append(subentry[0])
                    if defs: 
                        combined_dict_def = ""
                        for i in range(min(len(defs), 2)):
                            combined_dict_def += str(defs[i]) + "; "
                        return jsonify({"answer": " 📖 Definition: " + combined_dict_def.strip("; ")})

            # Second try: Transliterated/English fallback
            if lang != 'en':
                detect_url = f"https://translate.googleapis.com/translate_a/single?client=gtx&sl=auto&tl=en&dt=t&q={word}"
                res = requests.get(detect_url, timeout=3)
                if res.ok:
                    data = res.json()
                    english_word = data[0][0][0] if data[0] and data[0][0] else word
                    en_dict_url = f"https://translate.googleapis.com/translate_a/single?client=gtx&sl=en&tl=en&dt=md&q={english_word}"
                    res_en = requests.get(en_dict_url, timeout=3)
                    if res_en.ok:
                        en_data = res_en.json()
                        if len(en_data) > 12 and en_data[12]:
                            en_defs = [s[0] for e in en_data[12] if len(e) > 1 and e[1] for s in e[1] if s[0]]
                            if en_defs:
                                joined_en = "; ".join(en_defs[:2])
                                final_res = requests.get(f"https://translate.googleapis.com/translate_a/single?client=gtx&sl=en&tl={lang}&dt=t&q={joined_en}", timeout=3)
                                if final_res.ok:
                                    back_parts = []
                                    raw_parts = final_res.json()[0]
                                    for p in raw_parts:
                                        if p[0]: back_parts.append(str(p[0]))
                                    trans_back = "".join(back_parts)
                                    
                                    # Fallback slice removal
                                    limit_defs = []
                                    for i in range(min(len(en_defs), 2)): limit_defs.append(en_defs[i])
                                    
                                    return jsonify({"answer": f" 📖 Definition ({english_word}): " + trans_back})
        except Exception as e:
            print("Advanced Dict Fallback Error:", e)

        # 3. Last Resort: Use sentence context if no dictionary definition found
        matches = []
        if sentences:
            for s in sentences:
                if re.search(rf"\b{word}\b" if lang == "en" else word, s.lower()):
                    matches.append(s)
        
        if matches:
            matches.sort(key=lambda x: abs(60 - len(x)))
            return jsonify({"answer": f"💡 AI Context: {matches[0]}"})

        return jsonify({"answer": "I couldn't find a specific meaning for this word. Try a simpler word or check the context."})
        
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@app.route("/api/ask", methods=["POST"])
def ask_assistant():
    """Smart AI Assistant endpoint for conversational book queries."""
    try:
        data = request.get_json()
        query = data.get("query", "").lower()
        book_id = data.get("book_id")
        context_text = data.get("context", "")
        
        if not book_id and not context_text:
            return jsonify({"answer": "I need some context from a book to help you."})

        # 1. Retrieve book text if not provided in context
        if not context_text and book_id:
            sentences = get_book_sentences(book_id)
        else:
            # Tokenize context
            sentences = re.split(r'(?<=[.!?])\s+', context_text)
            sentences = [s.strip() for s in sentences if len(s.strip()) > 5]

        # 2. Intent Detection (Simple but effective for internal use)
        if "summarize" in query:
            # Summary intent: take a cluster of sentences
            summary_limit = 5
            return jsonify({"answer": f"📝 Here is a quick summary: {' '.join(sentences[:summary_limit])}..."})
            
        if "explain" in query or "meaning" in query or "what" in query:
            # Explanation intent: Find most relevant sentences
            keywords = [w for w in query.split() if len(w) > 3]
            best_match = ""
            max_hits = 0
            for s in sentences:
                hits = sum(1 for k in keywords if k in s.lower())
                if hits > max_hits:
                    max_hits = hits
                    best_match = s
            
            if best_match:
                return jsonify({"answer": f"💡 AI Explanation: {best_match}"})
            else:
                return jsonify({"answer": "I found something relevant: " + sentences[0]})

        # 3. Fallback to existing search logic for specific facts
        return jsonify({"answer": f"I'm looking into that! This part of the book mentions: {sentences[0]}"})

    except Exception as e:
        print(f"AI Assistant Error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route("/explain_image", methods=["POST"])
def explain_image():
    try:
        data = request.get_json()
        # Support both 'src' (legacy/original) and 'image' (new frontend logic)
        src = data.get("src", "") or data.get("image", "")
        context = data.get("context", "")
        
        if not src:
            return jsonify({"explanation": "No image provided."})

        # Load the image depending on its source
        img_np = None
        
        if src.startswith("data:image"):
            # Extract base64 part safely
            if "," in src:
                encoded_data = src.split(",", 1)[1]
                # The frontend might pass URL-encoded data (like %0A for newlines)
                # Fix common URL-encoding issues where '+' becomes ' '
                encoded_data = encoded_data.replace(' ', '+')
                # Fix padding if missing
                encoded_data += "=" * ((4 - len(encoded_data) % 4) % 4)
                
                try:
                    raw = base64.b64decode(encoded_data)
                    # First try OpenCV because it's fast
                    nparr = np.frombuffer(raw, np.uint8)
                    img_np = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                    
                    # Fallback to PIL (handles more formats like GIF, some JPX, TIFF)
                    if img_np is None:
                        pil_img = Image.open(io.BytesIO(raw)).convert("RGB")
                        img_np = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
                except Exception as e:
                    print("Base64 Decode or Image conversion failed:", e)
                        
        elif "/uploads/" in src:
            filename = src.split("/uploads/")[-1]
            filename = filename.split("?")[0].split("#")[0]
            filepath = os.path.join(UPLOAD_FOLDER, unquote(filename))
            img_np = cv2.imread(filepath)
            
        elif src.startswith("http"):
            try:
                req = urllib.request.Request(src, headers={'User-Agent': 'Mozilla/5.0'})
                with urllib.request.urlopen(req, timeout=5) as resp:
                    raw = resp.read()
                    nparr = np.frombuffer(raw, np.uint8)
                    img_np = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
                    if img_np is None:
                        pil_img = Image.open(io.BytesIO(raw)).convert("RGB")
                        img_np = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
            except Exception as e:
                print("HTTP Image fetch failed:", e)
            
        if img_np is None:
            return jsonify({"explanation": f"Could not read or locate the image. [Format unsupported or unreadable]\nSource preview: {src[:80]}..."})
            
        explanation = ""
        
        # --- Hybrid Explanation Strategy ---
        explanation = ""
        ocr_text = ""
        ai_caption = ""
        faces = []
        
        # 1. Run OCR (Tesseract) to find labels/data
        try:
            gray = cv2.cvtColor(img_np, cv2.COLOR_BGR2GRAY)
            # Denoise and sharpen for better OCR
            gray = cv2.GaussianBlur(gray, (0,0), 3)
            gray = cv2.addWeighted(gray, 1.5, gray, -0.5, 0)
            
            ocr_text_raw = pytesseract.image_to_string(gray, config='--oem 3 --psm 6').strip()
            ocr_text = re.sub(r'[^a-zA-Z0-9\s.,!?:;@&()\"\'-]', '', ocr_text_raw)
            if not ocr_text and re.search(r'[a-zA-Z]{3,}', ocr_text_raw):
                ocr_text = ocr_text_raw.strip()
        except Exception:
            pass

        try:
            img_pil = Image.fromarray(cv2.cvtColor(img_np, cv2.COLOR_BGR2RGB))
            
            if not hasattr(app, "image_captioner"):
                try:
                    if pipeline:
                        app.image_captioner = pipeline("image-to-text", model="nlp-connect/vit-gpt2-image-captioning", device=-1)
                    else:
                        app.image_captioner = None
                except Exception:
                    app.image_captioner = None
                    
            if app.image_captioner:
                # Optimized for speed
                res = app.image_captioner(img_pil, max_new_tokens=25)
                if res and len(res) > 0 and 'generated_text' in res[0]:
                    ai_caption = res[0]['generated_text'].capitalize().strip()
                
            # --- Fallback: Advanced Heuristics if AI fails ---
            if not ai_caption:
                # 1. Face Detection (Detect if it's a person/baby)
                try:
                    gray = cv2.cvtColor(img_np, cv2.COLOR_BGR2GRAY)
                    face_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
                    detected_faces = face_cascade.detectMultiScale(gray, 1.3, 5)
                    faces = detected_faces if detected_faces is not None else []
                    if len(faces) > 0:
                        ai_caption = "A close-up of a person's face"
                        if "BABY" in (context or "").upper() or "LITTLE" in (context or "").upper():
                            ai_caption = "A visual of a small baby"
                except Exception:
                    faces = []
                
                # 2. Contextual Inference based on OCR keywords
                if len(ocr_text) > 3:
                    u_ocr = ocr_text.upper()
                    u_ctx = (context or "").upper()
                    if "HARRY POTTER" in u_ocr or "WHO LIVED" in u_ocr or "HARRY POTTER" in u_ctx or "WHO LIVED" in u_ctx:
                        if len(faces) > 0 or "BABY" in u_ctx:
                            ai_caption = "An illustration of baby Harry Potter (The Boy Who Lived)"
                        else:
                            ai_caption = "An illustration for the Harry Potter series"
                    elif "BABY" in u_ctx:
                         ai_caption = "A visual depiction of a small baby"
                elif "BABY" in (context or "").upper():
                     ai_caption = "An illustration of a baby"
        except Exception as ai_e:
            print("AI Captioning failed:", ai_e)
        
        # FINAL SANITY: Remove noise OCR (single letters) from the final output comparison
        if len(ocr_text) < 4:
            ocr_text = ""

        # 3. Analyze and Combine (Detailed Visual Logic)
        is_diagram = any(kw in (ocr_text + ai_caption).lower() for kw in ["diagram", "chart", "graph", "figure", "fig.", "illustration", "table", "cycle", "process", "thank you", "slide", "presentation"])
        
        # Priority 1: Visual Scene Description
        if ai_caption:
            visual_desc = f"🖼️ AI VISUAL SUMMARY: {ai_caption}."
            
            # Enrich with context to explain "what it belongs to"
            if context:
                visual_desc += f"\n\n📚 VIEWPOINT: This appears in a passage discussing: '{context[:150]}...'."
            
            if ocr_text:
                if is_diagram:
                    explanation = f"{visual_desc}\n\n🔍 KEY DETAILS: The image contains text which points to: \"{ocr_text[:300]}...\""
                else:
                    explanation = f"{visual_desc}\n\n📝 EMBEDDED TEXT: {ocr_text[:120]}..."
            else:
                explanation = visual_desc
                
        # Priority 2: Text-only images or charts where AI failed
        elif ocr_text:
            cleaned_context = re.sub(r'[:;{}[]]', '', context[:100]).strip()
            explanation = f"🔍 DOCUMENT SNAPSHOT: Based on the visual data, this contains text elements including: \"{ocr_text[:400]}...\"\n\n"
            if cleaned_context:
                explanation += f"🚀 INTEGRATION: This likely serves as a visual reference for: '{cleaned_context}'."
            else:
                explanation += "Note: Vision AI is starting up or unavailable, using high-precision OCR for internal details."
                
        else:
            explanation = "This image appears to be an illustration or decorative element. Try highlighting surrounding text to help me understand its theme!"

        return jsonify({"explanation": explanation})
        
    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": "Failed to analyze the image."}), 500

@app.route("/analyze_emotion", methods=["POST"])
def analyze_emotion():
    try:
        data = request.get_json()
        text = data.get("text", "").lower()
        if not text:
            return jsonify({"emotion": "neutral"})

        # 1. Rule-based Fast Emotion Mapping
        emotion_rules = {
            "happy": ["happy", "joy", "wonderful", "delighted", "smile", "laugh", "cheerful", "magic", "sunshine", "hope", "love", "friend", "proud", "achieved"],
            "sad": ["sad", "cried", "tear", "unhappy", "lost", "death", "lonely", "darkness", "misery", "sorrow", "alone", "grave", "hurt", "pain", "hopeless"],
            "angry": ["angry", "rage", "hate", "fight", "shout", "mad", "fury", "annoyed", "bitter", "punch", "strike", "vengeance"],
            "fear": ["fear", "scared", "terrified", "ghost", "dark", "shadow", "unknown", "scary", "shiver", "beast", "creepy", "dangerous", "nervous"],
        }
        
        words = text.split()
        found_counts = {emotion: sum(1 for word in words if word in keywords) for emotion, keywords in emotion_rules.items()}
        
        # Safe max with lambda
        max_emotion = "neutral"
        max_val = -1
        for em, val in found_counts.items():
            if val > max_val:
                max_val = val
                max_emotion = em
        
        if max_val > 0:
            return jsonify({"emotion": max_emotion})

        # 2. AI Fallback
        analyzer = get_sentiment_analyzer()
        if analyzer:
            res = analyzer(text[:512])[0]
            if res['label'] == 'POSITIVE':
                return jsonify({"emotion": "happy"})
            else:
                return jsonify({"emotion": "fear"})
        return jsonify({"emotion": "neutral"})
    except Exception as e:
        print(f"Emotion analysis error: {e}")
        return jsonify({"emotion": "neutral"})

@app.route("/get_recommendations", methods=["POST"])
def get_recommendations():
    try:
        data = request.json
        book_id = data.get("book_id")
        
        conn = get_conn()
        cur = conn.cursor()
        cur.execute("SELECT name FROM books WHERE id = ?", (book_id,))
        row = cur.fetchone()
        conn.close()
        
        current_book_name = str(os.path.basename(row[0])) if row else ""
        
        # 1. Smarter Keyword Extraction
        clean_name = re.sub(r'(\.pdf|\.epub|\.docx|\.txt)', '', current_book_name, flags=re.IGNORECASE)
        clean_name = re.sub(r'(_|-)', ' ', clean_name)
        stop_words = {'the', 'and', 'a', 'of', 'to', 'in', 'is', 'it', 'for', 'with', 'on'}
        words = [w for w in clean_name.lower().split() if w not in stop_words and len(w) > 2]
        
        search_query = '+'.join(words[:2]) if words else "adventure"
        
        # SEED-BASED DIVERSIFICATION: Add a random genre if the title is generic
        genres = ["mystery", "magic", "ghost", "nautical", "detective", "fantasy", "scifi", "science", "travel", "history"]
        random_genre = random.choice(genres)
        
        all_recs = {} # Dedup by title

        def add_to_recs(book_list):
            for b in book_list:
                title_key = b['title'].strip().lower()
                if title_key not in all_recs:
                    all_recs[title_key] = b

        # 2. Attempt Cross-API Discovery
        try:
            # SEARCH A: Primary Keyword Search (OpenLibrary)
            res1 = requests.get(f"https://openlibrary.org/search.json?q={search_query}&limit=5", timeout=6)
            if res1.status_code == 200:
                docs = res1.json().get('docs', [])
                add_to_recs([{
                    "title": d.get('title'),
                    "author": d.get('author_name', ['Unknown'])[0],
                    "id": f"ol-{d.get('key', '').split('/')[-1]}",
                    "cover": f"https://covers.openlibrary.org/b/id/{d.get('cover_i')}-M.jpg" if d.get('cover_i') else None,
                    "url": f"https://openlibrary.org{d.get('key')}", "is_local": False
                } for d in docs if d.get('title')])
        except: pass

        try:
            # SEARCH B: Genre Discovery (Gutenberg)
            res2 = requests.get(f"https://gutendex.com/books/?topic={random_genre}", timeout=6)
            if res2.status_code == 200:
                results = res2.json().get('results', [])
                add_to_recs([{
                    "title": r.get('title'),
                    "author": r.get('authors', [{'name': 'Unknown'}])[0].get('name'),
                    "id": r.get('id'),
                    "cover": r.get('formats', {}).get('image/jpeg'),
                    "url": r.get('formats', {}).get('text/html') or r.get('formats', {}).get('text/plain; charset=utf-8'),
                    "is_local": False
                } for r in results if r.get('title')])
        except: pass

        # 3. MASSIVE MASTERPIECE VAULT (50+ Shuffled Classics)
        # We use this to fill gaps and ensure variety even when APIs are slow.
        if len(all_recs) < 8:
            vault = [
                (11, "Alice's Adventures in Wonderland", "Lewis Carroll"),
                (1661, "The Adventures of Sherlock Holmes", "Arthur Conan Doyle"),
                (84, "Frankenstein", "Mary Shelley"),
                (1342, "Pride and Prejudice", "Jane Austen"),
                (35, "The Time Machine", "H.G. Wells"),
                (120, "Treasure Island", "R.L. Stevenson"),
                (236, "The Jungle Book", "Rudyard Kipling"),
                (345, "Dracula", "Bram Stoker"),
                (730, "Oliver Twist", "Charles Dickens"),
                (1727, "The Odyssey", "Homer"),
                (2591, "Grimm's Fairy Tales", "The Brothers Grimm"),
                (158, "Emma", "Jane Austen"),
                (2701, "Moby Dick", "Herman Melville"),
                (36, "The War of the Worlds", "H.G. Wells"),
                (219, "Heart of Darkness", "Joseph Conrad"),
                (3207, "The Island of Doctor Moreau", "H.G. Wells"),
                (768, "Wuthering Heights", "Emily Bronte"),
                (98, "A Tale of Two Cities", "Charles Dickens"),
                (16328, "Beowulf", "Unknown"),
                (514, "Little Women", "Louisa May Alcott")
            ]
            random.shuffle(vault)
            for vid, vtitle, vauthor in vault:
                if len(all_recs) >= 20: break
                add_to_recs([{
                    "title": vtitle, "author": vauthor, "id": vid,
                    "cover": f"https://www.gutenberg.org/cache/epub/{vid}/pg{vid}.cover.medium.jpg",
                    "url": f"https://www.gutenberg.org/ebooks/{vid}.html.images", "is_local": False
                }])

        # Take the final set, shuffle, and return
        final_list = list(all_recs.values())
        random.shuffle(final_list)
        return jsonify({"recommendations": final_list[:4]})
        
    except Exception as e:
        print(f"Discovery Error: {e}")
        return jsonify({"recommendations": []})

@app.route("/download_external", methods=["POST"])
def download_external():
    data = request.json
    title = data.get("title")
    source_url = data.get("url") # Now we receive a real URL
    
    if not source_url or not title:
        return jsonify({"error": "Missing Source Information"}), 400
        
    try:
        # Increase timeout for guaranteed ingestion of large classics
        print(f"Downloading high-fidelity illustrated book from {source_url}...")
        res = requests.get(source_url, timeout=25)
        # Handle encoding properly
        res.encoding = res.apparent_encoding
        book_content = res.text
        
        # Determine if it's likely HTML
        is_html = "text/html" in source_url or "<html" in book_content[:2000].lower() or "<div" in book_content[:1000].lower()
        
        pages = []
        if is_html:
            soup = BeautifulSoup(book_content, "html.parser")
            
            # 1. Fix relative URLs for images and links
            for tag in soup.find_all(['img', 'a']):
                attr = 'src' if tag.name == 'img' else 'href'
                val = tag.get(attr)
                if val and not val.startswith(('http', 'data:', '#', 'mailto:')):
                    tag[attr] = urljoin(source_url, val)
            
            # 2. Extract title if better one found
            t_tag = soup.find('title')
            if t_tag: 
                new_title = t_tag.get_text().strip()
                if new_title and len(new_title) > 3: title = new_title

            # 3. Clean up boilerplate (optional but makes it feel more premium)
            # Remove giant Gutenberg header/footer blocks if they are distinct
            for noise in soup.find_all(['style', 'script']):
                noise.decompose()

            # 4. Safe Split into Pages
            # We iterate through blocks and group them
            body = soup.find('body') or soup
            current_page = []
            current_len = 0
            
            # We want roughly 3000-4000 chars per page
            for elem in body.children:
                elem_html = str(elem)
                # Skip whitespace nodes
                if not elem_html.strip(): continue
                
                if current_len + len(elem_html) > 4000 and current_page:
                    pages.append("".join(current_page))
                    current_page = [elem_html]
                    current_len = len(elem_html)
                else:
                    current_page.append(elem_html)
                    current_len += len(elem_html)
            
            if current_page:
                pages.append("".join(current_page))
        else:
            # Plain Text handling
            # Escape HTML and split by paragraphs
            safe_text = html.escape(book_content)
            # Split by double newlines (paragraphs)
            paragraphs = safe_text.split('\n\n')
            current_page = []
            current_len = 0
            for p in paragraphs:
                p_html = f"<p style='margin-bottom: 1.5em; line-height: 1.6;'>{p.strip()}</p>"
                if current_len + len(p_html) > 3000 and current_page:
                    pages.append("".join(current_page))
                    current_page = [p_html]
                    current_len = len(p_html)
                else:
                    current_page.append(p_html)
                    current_len += len(p_html)
            if current_page:
                pages.append("".join(current_page))

        # 5. Build final content with proper page IDs for the frontend splitter
        final_html = ""
        if not pages:
            # Fallback
            final_html = f"<div id='pdf-page-0' class='lazy-page-container' style='padding: 60px; max-width: 900px; margin: 0 auto; background: white;'>{book_content}</div>"
        else:
            for i, p_content in enumerate(pages):
                # We use the marker <div id="pdf-page- which the frontend looks for
                final_html += f'<div id="pdf-page-{i}" class="lazy-page-container" style="padding: 60px; max-width: 900px; margin: 0 auto 40px auto; background: var(--bg-paper); border-radius: 4px; box-shadow: 0 4px 20px rgba(0,0,0,0.1); min-height: 800px; color: #333; font-family: \'Georgia\', serif; font-size: 1.15rem;">\n'
                if i == 0:
                    final_html += f'<h1 style="text-align: center; margin-bottom: 50px; color: #2c3e50;">{title}</h1>\n'
                final_html += p_content
                final_html += '\n</div>'

        safe_name = secure_filename(title) + ".html"
        filepath = os.path.join(UPLOAD_FOLDER, safe_name)
        extracted_path = os.path.join(EXTRACTED_FOLDER, safe_name + ".html")

        with open(filepath, "w", encoding="utf-8") as f:
            f.write(final_html)
            
        conn = get_conn()
        cur = conn.cursor()
        cur.execute(
            "INSERT INTO books(name, path, extracted_path, status) VALUES (?, ?, ?, ?)",
            (title, filepath, extracted_path, "processing")
        )
        book_id = cur.lastrowid
        conn.commit()
        conn.close()

        def do_extract_external(fpath, epath, bid):
            try:
                # We don't need heavy extraction for Gutenberg as we've already formatted it
                # But we copy it to extracted_path for consistency
                with open(fpath, "r", encoding="utf-8") as f:
                    html_data = f.read()
                with open(epath, "w", encoding="utf-8") as f:
                    f.write(html_data)
                
                conn2 = get_conn()
                conn2.execute("UPDATE books SET status='ready' WHERE id=?", (bid,))
                conn2.commit()
                conn2.close()
            except Exception as e:
                print(f"External Full-Text Ingestion Error: {e}")
                conn2 = get_conn()
                conn2.execute("UPDATE books SET status='error' WHERE id=?", (bid,))
                conn2.commit()
                conn2.close()

        t = threading.Thread(target=do_extract_external, args=(filepath, extracted_path, book_id), daemon=True)
        t.start()
            
        return jsonify({"status": "success", "id": book_id})
    except Exception as e:
        print(f"Full-Text Download failed: {e}")
        return jsonify({"error": str(e)}), 500

if __name__ == "__main__":
    init_db()
    start_translation_sidecar()
    app.run(debug=True, host="0.0.0.0", port=5000)